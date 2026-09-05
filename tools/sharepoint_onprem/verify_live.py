"""Opt-in real-server smoke test. Creates only a uniquely named QA folder.

Run from backend with its venv:
  python ../tools/sharepoint_onprem/verify_live.py --site-url https://portal/sites/lab --username 'DOMAIN\\reader'
Use --seed with an account allowed to upload into Documents. Password is read
from SHAREPOINT_TEST_PASSWORD or a hidden prompt, never command-line arguments.
"""
import argparse
import io
import json
import os
import sys
import uuid
import zipfile
from getpass import getpass
from pathlib import Path
from urllib.parse import quote

sys.path.insert(0, str(Path(__file__).resolve().parents[2] / "backend"))
import pandas as pd
import requests
from requests_ntlm import HttpNtlmAuth
from pypdf import PdfWriter
from pypdf.generic import DictionaryObject, NameObject, DecodedStreamObject
from pptx import Presentation
from app.data_sources.clients.sharepoint_onprem_client import SharepointOnpremClient
from app.data_sources.clients._file_source_common import GlobScopeError


def fixtures():
    files = {
        "sales.csv": b"region,revenue\nNorth,1200\nSouth,800\nWest,1500\n",
        "policy.txt": b"BOW connector QA policy. The travel reimbursement limit is 450 euros per trip. Approval is required from the Finance team.",
        "metadata.json": b'{"project":"Orion","owner":"Finance","budget":3500}',
        "O'Brien #100%.txt": b"Special filename verification. The project code is ORION-731.",
    }
    buf = io.BytesIO()
    with pd.ExcelWriter(buf) as writer:
        pd.DataFrame({"region": ["North", "South", "West"], "revenue": [1200, 800, 1500]}).to_excel(writer, sheet_name="Sales", index=False)
        pd.DataFrame({"target": [4000]}).to_excel(writer, sheet_name="Targets", index=False)
    files["sales.xlsx"] = buf.getvalue()
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as z:
        z.writestr("[Content_Types].xml", '<?xml version="1.0"?><Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types"><Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/><Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/></Types>')
        z.writestr("_rels/.rels", '<?xml version="1.0"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/></Relationships>')
        z.writestr("word/document.xml", '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"><w:body><w:p><w:r><w:t>Orion launch briefing. The launch owner is Finance and the budget is 3500 euros.</w:t></w:r></w:p></w:body></w:document>')
    files["briefing.docx"] = buf.getvalue()
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[1])
    slide.shapes.title.text = "Orion status"
    slide.placeholders[1].text = "Launch readiness is green. Finance owns the launch. Revenue is 3500 euros."
    buf = io.BytesIO(); pres.save(buf); files["status.pptx"] = buf.getvalue()
    writer = PdfWriter()
    font = DictionaryObject({NameObject("/Type"): NameObject("/Font"), NameObject("/Subtype"): NameObject("/Type1"), NameObject("/BaseFont"): NameObject("/Helvetica")})
    for text in ["Orion monthly report. Revenue is 3500 euros across three regions.", "The travel reimbursement limit is 450 euros per trip. Finance approval is required."]:
        page = writer.add_blank_page(width=612, height=792)
        page[NameObject("/Resources")] = DictionaryObject({NameObject("/Font"): DictionaryObject({NameObject("/F1"): writer._add_object(font)})})
        stream = DecodedStreamObject()
        stream.set_data(f"BT /F1 12 Tf 50 720 Td ({text}) Tj ET".encode())
        page[NameObject("/Contents")] = writer._add_object(stream)
    buf = io.BytesIO(); writer.write(buf); files["report.pdf"] = buf.getvalue()
    return files


def main():
    p = argparse.ArgumentParser(description=__doc__)
    p.add_argument("--site-url", required=True)
    p.add_argument("--username", required=True)
    p.add_argument("--allow-http", action="store_true")
    p.add_argument("--seed", action="store_true")
    p.add_argument("--folder", default="")
    args = p.parse_args()
    password = os.environ.get("SHAREPOINT_TEST_PASSWORD") or getpass("SharePoint lab password: ")
    folder = args.folder
    # Validate the transport policy before any optional fixture upload.
    SharepointOnpremClient(site_url=args.site_url, username=args.username,
                          password=password, allow_http=args.allow_http)
    fixture_bytes = fixtures()
    if args.seed:
        session = requests.Session(); session.trust_env = False
        session.auth = HttpNtlmAuth(args.username, password)
        session.headers["Accept"] = "application/json;odata=nometadata"
        digest = session.post(args.site_url + "/_api/contextinfo", timeout=30)
        digest.raise_for_status()
        session.headers["X-RequestDigest"] = digest.json()["FormDigestValue"]
        root = session.get(args.site_url + "/_api/web/lists/getbytitle('Documents')/RootFolder", timeout=30).json()["ServerRelativeUrl"]
        folder = "BOW-QA-" + uuid.uuid4().hex[:8]
        path = root + "/" + folder
        r = session.post(args.site_url + "/_api/web/folders", json={"ServerRelativeUrl": path}, timeout=30); r.raise_for_status()
        for name, blob in fixture_bytes.items():
            url = args.site_url + "/_api/web/GetFolderByServerRelativePath(decodedurl='" + quote(path.replace("'", "''"), safe="/") + "')/Files/AddUsingPath(decodedurl='" + quote(name.replace("'", "''"), safe="") + "',overwrite=false)"
            r = session.post(url, data=blob, headers={"Content-Type": "application/octet-stream"}, timeout=45)
            r.raise_for_status()
        print("Created QA fixture folder:", folder, flush=True)
    c = SharepointOnpremClient(site_url=args.site_url, username=args.username, password=password,
                              allow_http=args.allow_http, drive_name="Documents", folder_path=folder, recursive=True)
    probe = c.test_connection()
    assert probe["success"], probe["message"]
    print("PASS authenticated site and scoped library probe", flush=True)
    files = c.list_files()
    print("PASS list files:", len(files), flush=True)
    assert len(files) == len(fixture_bytes) if args.seed else len(files) > 0
    assert len(c.get_schemas()) == len(files)
    for f in files:
        result = c.read_file(f["id"])
        raw, name, mime = c.read_raw_bytes(f["id"])
        assert raw and name == f["name"]
        if args.seed:
            assert raw == fixture_bytes[name]
        print("PASS read + original bytes:", name, type(result).__name__, flush=True)
    if args.seed or folder.startswith("BOW-QA-"):
        assert c.read_file("sales.xlsx", sheet="Targets")["target"].iloc[0] == 4000
        assert pd.to_numeric(c.read_file("sales.csv")["revenue"]).sum() == 3500
        page = c.read_file("report.pdf", page_range=(2, 2))
        assert page["pages_total"] == 2 and "450" in page["text"]
        assert "Finance" in c.read_file("briefing.docx")
        assert "3500" in c.read_file("status.pptx")
        assert c.search_files("policy")[0]["name"] == "policy.txt"
        print("PASS Excel sheets, CSV totals, PDF page 2, DOCX/PPTX text, live filename search", flush=True)
        scoped = SharepointOnpremClient(site_url=args.site_url, username=args.username, password=password,
                    allow_http=args.allow_http, drive_name="Documents", folder_path=folder, include_globs="*.csv")
        assert len(scoped.list_files()) == 1
        try:
            scoped.read_file(next(f["id"] for f in files if f["name"] == "policy.txt"))
        except GlobScopeError:
            print("PASS off-glob read denied", flush=True)
        else:
            raise AssertionError("Scope bypass")
    print(json.dumps({"site_url": args.site_url, "folder": folder, "files": [f["name"] for f in files]}), flush=True)


if __name__ == "__main__":
    main()
