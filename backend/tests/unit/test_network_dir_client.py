"""Unit tests for the NetworkDirClient (network_dir data source).

Exercises the filesystem primitives — list / search (filename + content) /
read / write / copy — plus the security invariants (root confinement, path
traversal rejection, read-only enforcement) against a generated temp tree.
No DB or network.
"""
from __future__ import annotations

from pathlib import Path

import pandas as pd
import pytest

from app.data_sources.clients.base import Capability
from app.data_sources.clients.network_dir_client import NetworkDirClient


@pytest.fixture()
def tree(tmp_path: Path) -> Path:
    """A small deterministic directory tree."""
    (tmp_path / "contracts").mkdir()
    (tmp_path / "images").mkdir()
    (tmp_path / "contracts" / "acme_2025.csv").write_text("clause,penalty\npayment,100\n")
    (tmp_path / "contracts" / "globex_2024.csv").write_text("clause,penalty\nrenewal,50\n")
    (tmp_path / "contracts" / "acme_summary.md").write_text(
        "# Acme Master Agreement\nStatus: active. Vendor Umbrella referenced here.\n"
    )
    (tmp_path / "notes.txt").write_text("misc notes about invoices\n")
    # a tiny binary (fake png header)
    (tmp_path / "images" / "logo.png").write_bytes(b"\x89PNG\r\n\x1a\n" + b"\x00" * 32)
    return tmp_path


# ---------------------------------------------------------------- registry


class TestRegistry:
    def test_registered_and_file_shaped(self):
        from app.schemas.data_source_registry import REGISTRY, resolve_client_class

        entry = REGISTRY["network_dir"]
        assert entry.is_connection is True
        assert entry.data_shape == "files"
        assert entry.catalog_ownership == "shared"
        assert resolve_client_class("network_dir") is NetworkDirClient

    def test_type_is_a_file_source(self):
        from app.ai.tools.implementations._file_tool_common import FILE_SOURCE_TYPES

        assert "network_dir" in FILE_SOURCE_TYPES


# ------------------------------------------------------------ capabilities


class TestCapabilities:
    def test_class_advertises_write(self):
        assert Capability.WRITE_FILE in NetworkDirClient.capabilities

    def test_readonly_instance_drops_write(self, tree):
        c = NetworkDirClient(root_path=str(tree), writable=False)
        assert Capability.WRITE_FILE not in c.capabilities
        assert Capability.LIST_FILES in c.capabilities

    def test_writable_instance_keeps_write(self, tree):
        c = NetworkDirClient(root_path=str(tree), writable=True)
        assert Capability.WRITE_FILE in c.capabilities


# --------------------------------------------------------------- reads


class TestReadOnly:
    def test_test_connection(self, tree):
        c = NetworkDirClient(root_path=str(tree))
        res = c.test_connection()
        assert res["success"] is True

    def test_list_files_recursive(self, tree):
        c = NetworkDirClient(root_path=str(tree), recursive=True)
        files = c.list_files()
        ids = {f["id"] for f in files}
        assert "contracts/acme_2025.csv" in ids
        assert "images/logo.png" in ids
        # ids are stable relative POSIX paths
        assert all("\\" not in f["id"] for f in files)

    def test_list_files_non_recursive(self, tree):
        c = NetworkDirClient(root_path=str(tree), recursive=False)
        ids = {f["id"] for f in c.list_files()}
        assert "notes.txt" in ids
        assert "contracts/acme_2025.csv" not in ids  # subfolder excluded

    def test_extension_filter(self, tree):
        c = NetworkDirClient(root_path=str(tree), allowed_extensions="csv")
        exts = {f["name"].rsplit(".", 1)[-1] for f in c.list_files()}
        assert exts == {"csv"}

    def test_read_csv_returns_dataframe(self, tree):
        c = NetworkDirClient(root_path=str(tree))
        df = c.read_file("contracts/acme_2025.csv")
        assert isinstance(df, pd.DataFrame)
        assert list(df.columns) == ["clause", "penalty"]

    def test_read_text(self, tree):
        c = NetworkDirClient(root_path=str(tree))
        text = c.read_file("notes.txt")
        assert isinstance(text, str)
        assert "invoices" in text

    def test_read_binary_returns_bytes(self, tree):
        c = NetworkDirClient(root_path=str(tree))
        b = c.read_file("images/logo.png")
        assert isinstance(b, (bytes, bytearray))
        assert b[:4] == b"\x89PNG"

    def test_search_by_filename(self, tree):
        c = NetworkDirClient(root_path=str(tree))
        hits = c.search_files("acme", content=False)
        ids = {h["id"] for h in hits}
        assert "contracts/acme_2025.csv" in ids
        assert "contracts/globex_2024.csv" not in ids

    def test_search_by_content(self, tree):
        c = NetworkDirClient(root_path=str(tree))
        # "Umbrella" only appears INSIDE acme_summary.md, not in any filename.
        hits = c.search_files("Umbrella")
        ids = {h["id"] for h in hits}
        assert "contracts/acme_summary.md" in ids

    def test_prompt_schema_lists_files(self, tree):
        c = NetworkDirClient(root_path=str(tree))
        text = c.prompt_schema()
        assert "acme_2025.csv" in text


# --------------------------------------------------------------- security


class TestSecurity:
    def test_read_traversal_rejected(self, tree):
        c = NetworkDirClient(root_path=str(tree))
        with pytest.raises(ValueError, match="escapes"):
            c.read_file("../../etc/passwd")

    def test_absolute_escape_rejected(self, tree):
        c = NetworkDirClient(root_path=str(tree), writable=True)
        with pytest.raises(ValueError, match="escapes"):
            c.read_file("/etc/passwd")

    def test_write_on_readonly_rejected(self, tree):
        c = NetworkDirClient(root_path=str(tree), writable=False)
        with pytest.raises(ValueError, match="read-only"):
            c.write_file("x.txt", "hi")

    def test_write_traversal_rejected(self, tree):
        c = NetworkDirClient(root_path=str(tree), writable=True)
        with pytest.raises(ValueError, match="escapes"):
            c.write_file("../escape.txt", "x")

    def test_missing_root_errors(self):
        c = NetworkDirClient(root_path="/no/such/dir/here")
        assert c.test_connection()["success"] is False


# --------------------------------------------------------------- writes


class TestIndexingAndKeywords:
    def test_get_schemas_extracts_keywords(self, tree):
        c = NetworkDirClient(root_path=str(tree), index_content=True, max_keywords=50)
        tables = {t.name: t for t in c.get_schemas()}
        meta = tables["contracts/acme_2025.csv"].metadata_json["network_dir"]
        assert meta["indexed"] is True
        assert meta["content_hash"]
        # keywords include filename + content tokens
        kws = set(meta["keywords"])
        assert "acme" in kws
        assert "payment" in kws  # a cell value in the csv

    def test_index_content_off_stores_no_keywords(self, tree):
        c = NetworkDirClient(root_path=str(tree), index_content=False)
        meta = c.get_schemas()[0].metadata_json["network_dir"]
        assert "keywords" not in meta

    def test_keywords_capped(self, tmp_path):
        (tmp_path / "big.txt").write_text(" ".join(f"word{i}" for i in range(500)))
        c = NetworkDirClient(root_path=str(tmp_path), max_keywords=10)
        meta = c.get_schemas()[0].metadata_json["network_dir"]
        assert len(meta["keywords"]) <= 10

    def test_junk_files_skipped(self, tree):
        (tree / ".DS_Store").write_text("junk")
        (tree / "~$acme_2025.csv").write_text("lock stub")
        (tree / ".hidden").write_text("x")
        c = NetworkDirClient(root_path=str(tree))
        names = {f["name"] for f in c.list_files()}
        assert ".DS_Store" not in names
        assert "~$acme_2025.csv" not in names
        assert ".hidden" not in names

    def test_read_raw_bytes_returns_original(self, tree):
        c = NetworkDirClient(root_path=str(tree))
        data, name, mime = c.read_raw_bytes("images/logo.png")
        assert name == "logo.png"
        assert data[:4] == b"\x89PNG"
        assert mime == "image/png"

    def test_excel_sheet_names_and_cells_indexed(self, tmp_path):
        import pandas as pd
        path = tmp_path / "budget.xlsx"
        with pd.ExcelWriter(path) as xl:
            pd.DataFrame({"team": ["eng"], "count": [5]}).to_excel(xl, sheet_name="headcount", index=False)
        c = NetworkDirClient(root_path=str(tmp_path))
        # sheet name "headcount" is searchable even though it's not a cell value
        assert any(h["name"] == "budget.xlsx" for h in c.search_files("headcount"))
        # keyword index picks it up too
        kws = c.get_schemas()[0].metadata_json["network_dir"]["keywords"]
        assert "headcount" in kws


@pytest.fixture()
def doc_tree(tmp_path: Path) -> Path:
    """A directory with one pdf, one docx and one pptx, each containing a
    unique term that appears ONLY in that file's content (not its name)."""
    from scripts.gen_network_dir_fixtures import _make_docx, _make_pdf
    from pptx import Presentation
    from pptx.util import Inches

    (tmp_path / "docs").mkdir()
    _make_pdf(tmp_path / "docs" / "a.pdf", "Agreement",
              ["This PDF mentions zebracrossing as a secret term."])
    _make_docx(tmp_path / "docs" / "b.docx",
               ["A Word contract containing wombatclause internally."])
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2)).text_frame.text = (
        "Slide deck referencing quokkaterm in the body."
    )
    prs.save(tmp_path / "docs" / "c.pptx")
    return tmp_path


class TestDocumentExtraction:
    """pdf / docx / pptx must be readable and content-searchable, not opaque."""

    def test_read_pdf_returns_text(self, doc_tree):
        c = NetworkDirClient(root_path=str(doc_tree))
        out = c.read_file("docs/a.pdf")
        assert isinstance(out, str)
        assert "zebracrossing" in out

    def test_read_docx_returns_text(self, doc_tree):
        c = NetworkDirClient(root_path=str(doc_tree))
        out = c.read_file("docs/b.docx")
        assert isinstance(out, str)
        assert "wombatclause" in out

    def test_read_pptx_returns_text(self, doc_tree):
        c = NetworkDirClient(root_path=str(doc_tree))
        out = c.read_file("docs/c.pptx")
        assert isinstance(out, str)
        assert "quokkaterm" in out

    def test_search_matches_pdf_content(self, doc_tree):
        c = NetworkDirClient(root_path=str(doc_tree))
        ids = {h["id"] for h in c.search_files("zebracrossing")}
        assert "docs/a.pdf" in ids

    def test_search_matches_docx_content(self, doc_tree):
        c = NetworkDirClient(root_path=str(doc_tree))
        ids = {h["id"] for h in c.search_files("wombatclause")}
        assert "docs/b.docx" in ids

    def test_search_matches_pptx_content(self, doc_tree):
        c = NetworkDirClient(root_path=str(doc_tree))
        ids = {h["id"] for h in c.search_files("quokkaterm")}
        assert "docs/c.pptx" in ids

    def test_corrupt_document_does_not_break_search(self, doc_tree):
        # A file with a doc extension but garbage content must not raise;
        # extraction returns "" and search simply skips it.
        (doc_tree / "docs" / "broken.pdf").write_bytes(b"not a real pdf")
        c = NetworkDirClient(root_path=str(doc_tree))
        hits = c.search_files("zebracrossing")  # should still succeed
        assert any(h["id"] == "docs/a.pdf" for h in hits)


class TestWrites:
    def test_write_text(self, tree):
        c = NetworkDirClient(root_path=str(tree), writable=True)
        entry = c.write_file("out/summary.md", "# hi\n")
        assert entry["id"] == "out/summary.md"
        assert (tree / "out" / "summary.md").read_text() == "# hi\n"

    def test_overwrite_guard(self, tree):
        c = NetworkDirClient(root_path=str(tree), writable=True)
        c.write_file("dup.txt", "a")
        with pytest.raises(ValueError, match="already exists"):
            c.write_file("dup.txt", "b")
        # overwrite=True succeeds
        entry = c.write_file("dup.txt", "b", overwrite=True)
        assert entry["size"] == 1
        assert (tree / "dup.txt").read_text() == "b"

    def test_write_with_folder_arg(self, tree):
        c = NetworkDirClient(root_path=str(tree), writable=True)
        entry = c.write_file("acme.csv", "a,b\n1,2\n", folder_id="related")
        assert entry["id"] == "related/acme.csv"

    def test_write_extension_filter(self, tree):
        c = NetworkDirClient(root_path=str(tree), writable=True, allowed_extensions="csv,md")
        with pytest.raises(ValueError, match="not allowed"):
            c.write_file("bad.exe", "x")

    def test_copy_file_within_root(self, tree):
        c = NetworkDirClient(root_path=str(tree), writable=True)
        entry = c.copy_file("contracts/acme_2025.csv", "backup/acme_copy.csv")
        assert entry["id"] == "backup/acme_copy.csv"
        assert (tree / "backup" / "acme_copy.csv").exists()

    def test_write_bytes(self, tree):
        c = NetworkDirClient(root_path=str(tree), writable=True)
        entry = c.write_file("data/blob.bin", b"\x00\x01\x02")
        assert entry["size"] == 3
