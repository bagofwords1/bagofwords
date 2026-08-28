import asyncio
import base64
import json
import logging
import time
from pathlib import Path
from typing import AsyncIterator, Dict, Any, Type, List, Optional

import aiofiles
from pydantic import BaseModel
from sqlalchemy import select
from sqlalchemy.orm import selectinload

from app.ai.tools.base import Tool
from app.models.file import File
from app.models.report_file_association import report_file_association

logger = logging.getLogger(__name__)


from app.ai.tools.metadata import ToolMetadata
from app.ai.tools.schemas import (
    ToolEvent,
    ToolStartEvent,
    ToolProgressEvent,
    ToolEndEvent,
)
from app.ai.tools.schemas.create_artifact import CreateArtifactInput, CreateArtifactOutput
from app.ai.llm import LLM
from app.ai.llm.types import ImageInput, Message, TextDeltaEvent
from app.models.artifact import Artifact
from app.models.visualization import Visualization
from app.dependencies import async_session_maker
from app.services.thumbnail_service import ThumbnailService
from app.services.artifact_libs import get_inline_scripts
from app.ai.code_execution.pptx_executor import PptxCodeExecutor, PptxPreviewService
from sqlalchemy import desc
from app.ai.tools.implementations._sandbox_context import (
    SANDBOX_RUNTIME_PROMPT,
    ANON_PREVIEW_NOTE,
    build_identity_context,
)
from app.ai.tools.implementations._artifact_images import load_image_bytes
from app.ai.prompt_language import build_language_directive


class CreateArtifactTool(Tool):
    """Tool for generating React-based artifact code for dashboards.

    This tool generates standalone React/JSX code that renders visualizations
    using ECharts, styled with Tailwind CSS, and transpiled in-browser via Babel.

    The generated code runs in a sandboxed iframe and receives visualization
    data via window.ARTIFACT_DATA.
    """

    @property
    def metadata(self) -> ToolMetadata:
        return ToolMetadata(
            name="create_artifact",
            description=(
                "Create or fully rebuild artifacts (dashboards, pages, slide presentations) from visualizations. "
                "Use for: new dashboards, full redesigns, large layout changes, or when edit_artifact cannot handle the scope. "
                "Modes: 'page' for interactive dashboards with KPI cards, charts, and responsive grids; "
                "'slides' for presentation decks (exportable to PPTX). "
                "IMPORTANT: for 'page' mode visualization_ids are required - find them in previous create_data tool results "
                "shown as 'viz_id: <uuid>' in the conversation history. For 'slides' mode they are optional: "
                "a deck may include title, agenda and narrative slides that carry no chart. "
                "Do NOT ask the user for URLs or IDs - extract them from the conversation context. "
                "Only visualizations with successful step status are included."
            ),
            category="action",
            version="1.0.0",
            input_schema=CreateArtifactInput.model_json_schema(),
            output_schema=CreateArtifactOutput.model_json_schema(),
            max_retries=1,
            timeout_seconds=120,
            idempotent=False,
            required_permissions=[],
            is_active=True,
            tags=["artifact",  "dashboard", "slides"],
            allowed_modes=["chat"],
        )

    @property
    def input_model(self) -> Type[BaseModel]:
        return CreateArtifactInput

    @property
    def output_model(self) -> Type[BaseModel]:
        return CreateArtifactOutput

    # Path to the sandbox HTML file (relative to project root)
    # __file__ -> implementations -> tools -> ai -> app -> backend -> project_root
    SANDBOX_HTML_PATH = Path(__file__).parent.parent.parent.parent.parent.parent / "frontend" / "public" / "artifact-sandbox.html"

    async def _take_preview_screenshot(
        self,
        html_content: str,
    ) -> tuple[Optional[str], list[str]]:
        """Take a quick screenshot for planner reflection and capture JS errors.

        Returns (base64-encoded PNG string or None, list of JS error messages).
        """
        try:
            from playwright.async_api import async_playwright
        except ImportError:
            return None, []

        js_errors: list[str] = []

        try:
            import tempfile, os
            async with async_playwright() as p:
                # Optional executable override for deployments where the
                # Playwright-managed browser download is unavailable but a
                # compatible Chromium exists on disk.
                _exe = os.environ.get("BOW_CHROMIUM_EXECUTABLE") or None
                browser = await p.chromium.launch(headless=True, executable_path=_exe)
                page = await browser.new_page(viewport={"width": 1280, "height": 720})

                # Capture JS errors during render. Both channels matter:
                # pageerror catches thrown/uncaught exceptions (incl. Babel
                # transform errors), console 'error' catches failures that
                # libraries report without throwing (React render errors).
                def _on_pageerror(err):
                    if len(js_errors) < 10:
                        js_errors.append(str(err))

                def _on_console(msg):
                    try:
                        if msg.type == "error" and len(js_errors) < 10:
                            text = msg.text or ""
                            # Skip noise that isn't a code defect
                            if "favicon" in text.lower():
                                return
                            entry = f"[console.error] {text}"
                            if not any(text in e for e in js_errors):
                                js_errors.append(entry)
                    except Exception:
                        pass

                page.on("pageerror", _on_pageerror)
                page.on("console", _on_console)

                # Write HTML to a temp file and navigate via file:// URL.
                # This allows vendored scripts (e.g. Tailwind runtime) that use
                # document.write() to work correctly — document.write fails on
                # about:blank pages used by set_content().
                tmp = tempfile.NamedTemporaryFile(suffix=".html", delete=False, mode="w", encoding="utf-8")
                try:
                    tmp.write(html_content)
                    tmp.close()
                    await page.goto(f"file://{tmp.name}", wait_until="networkidle")

                    # Wait for React to mount and charts to render (short timeout)
                    try:
                        await page.wait_for_function(
                            "window.__ARTIFACT_RENDER_COMPLETE__ === true",
                            timeout=8000,
                        )
                    except Exception:
                        pass  # Take screenshot anyway — partial render is still useful

                    await asyncio.sleep(0.3)
                    screenshot_bytes = await page.screenshot(type="png", full_page=False)
                    await browser.close()
                    return base64.b64encode(screenshot_bytes).decode("utf-8"), js_errors
                finally:
                    os.unlink(tmp.name)
        except Exception as e:
            logger.warning(f"Preview screenshot failed: {e}")
            return None, js_errors

    async def _generate_thumbnail_background(
        self,
        artifact_id: str,
        html_content: str,
        mode: str = "page",
    ) -> None:
        """Generate thumbnail in background and update artifact.

        Runs independently with its own database session.
        """
        try:
            thumbnail_service = ThumbnailService()
            thumbnail_path = await thumbnail_service.generate_thumbnail(
                artifact_id=artifact_id,
                html_content=html_content,
                mode=mode,
            )
            if thumbnail_path:
                # Use a fresh database session for the background update
                async with async_session_maker() as db:
                    from sqlalchemy import update
                    from app.models.artifact import Artifact
                    stmt = update(Artifact).where(Artifact.id == artifact_id).values(thumbnail_path=thumbnail_path)
                    await db.execute(stmt)
                    await db.commit()
        except Exception as e:
            logger.warning(f"Failed to generate thumbnail for artifact {artifact_id}: {e}")

    async def _load_completion_images(
        self,
        db: Any,
        head_completion_id: Optional[str],
    ) -> List[ImageInput]:
        """Load images attached to the head completion as ImageInput objects.

        Args:
            db: Database session
            head_completion_id: The completion ID to load images for

        Returns:
            List of ImageInput objects ready for vision-capable LLM
        """
        if not head_completion_id:
            return []

        images: List[ImageInput] = []
        try:
            # Query files associated with this completion that are images
            result = await db.execute(
                select(File)
                .join(report_file_association, report_file_association.c.file_id == File.id)
                .where(report_file_association.c.completion_id == head_completion_id)
                .where(File.content_type.startswith("image/"))
            )
            image_files = result.scalars().all()

            for f in image_files:
                if not f.path:
                    continue
                try:
                    async with aiofiles.open(f.path, 'rb') as file:
                        content = await file.read()
                    images.append(ImageInput(
                        data=base64.b64encode(content).decode('utf-8'),
                        media_type=f.content_type or 'image/png',
                        source_type='base64'
                    ))
                except Exception as e:
                    logger.warning(f"Failed to load image file {f.id}: {e}")

        except Exception as e:
            logger.warning(f"Failed to query completion images: {e}")

        return images

    def _build_thumbnail_html(self, artifact_data: dict, code: str, mode: str = "page") -> str:
        """Build HTML for thumbnail generation in headless browser.

        Args:
            artifact_data: The data to inject as window.ARTIFACT_DATA
            code: The LLM-generated artifact code
            mode: 'page' for React dashboards, 'slides' for pure HTML presentations

        Returns:
            Complete HTML string ready for headless browser rendering
        """
        data_json = json.dumps(artifact_data, default=str)

        # Slides mode: pure HTML + Tailwind (no React/Babel)
        if mode == "slides":
            slides_scripts = get_inline_scripts(mode="slides")
            slides_template = """<!DOCTYPE html>
<html>
<head>
  <meta charset="UTF-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  __SLIDES_SCRIPTS__
  <style>
    html, body { height: 100%; margin: 0; padding: 0; }
    body { font-family: system-ui, -apple-system, sans-serif; }
    .slide { transition: opacity 0.3s ease-in-out; }
  </style>
</head>
<body class="bg-slate-900">
  <script>
    window.ARTIFACT_DATA = __ARTIFACT_DATA_JSON__;
    window.__ARTIFACT_RENDER_COMPLETE__ = false;
    setTimeout(function() {
      window.__ARTIFACT_RENDER_COMPLETE__ = true;
    }, 500);
  </script>

  __LLM_GENERATED_CODE__
</body>
</html>"""
            return slides_template.replace("__SLIDES_SCRIPTS__", slides_scripts).replace("__ARTIFACT_DATA_JSON__", data_json).replace("__LLM_GENERATED_CODE__", code)

        # Page mode: Build self-contained HTML mirroring ArtifactFrame.vue's approach.
        # get_inline_scripts("page") already includes all vendored libs + artifact-globals.js
        # so we only need to inject ARTIFACT_DATA, the LLM code, and render-complete detection.
        page_scripts = get_inline_scripts(mode="page")
        SC = '</' + 'script>'  # Avoid parser issues in this Python string too

        html = f"""<!DOCTYPE html>
<html>
<head>
  <meta charset="UTF-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  {page_scripts}
  <style>
    html, body, #root {{ height: 100%; margin: 0; padding: 0; }}
    body {{ font-family: system-ui, -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif; }}
  </style>
</head>
<body>
  <div id="root"></div>

  <script>
    window.ARTIFACT_DATA = {data_json};
    window.__ARTIFACT_RENDER_COMPLETE__ = false;
    window.__BOW_INFO = false;
  {SC}

  <script>
    // Params-runtime safety net for the headless validation render: the real
    // useParams/useParamOptions live in artifact-globals.js and need no host
    // to merely render, but if the loaded globals bundle predates them (stale
    // build output) correct artifact code would phantom-fail with
    // "useParams is not defined" and burn the repair loop. Stubs mirror the
    // no-params shape; they are defined ONLY when the globals didn't.
    if (typeof window.useParams !== 'function') {{
      window.useParams = function() {{
        return {{ declarations: [], values: {{}}, pending: {{}}, loading: false,
                 error: null, setParam: function() {{}}, setParams: function() {{}},
                 apply: function() {{}}, refresh: function() {{}},
                 getOptions: function() {{ return null; }} }};
      }};
    }}
    if (typeof window.useParamOptions !== 'function') {{
      window.useParamOptions = function() {{ return null; }};
    }}
    if (typeof window.useCurrentUser !== 'function') {{
      window.useCurrentUser = function() {{
        var d = window.ARTIFACT_DATA || {{}};
        return d.current_user || null;
      }};
    }}
  {SC}

  {code}

  <script>
    (function detectRenderComplete() {{
      var startTime = Date.now();
      var MAX_WAIT = 15000;
      function check() {{
        if (Date.now() - startTime > MAX_WAIT) {{
          window.__ARTIFACT_RENDER_COMPLETE__ = true;
          return;
        }}
        var root = document.getElementById('root');
        if (!root || root.children.length === 0) {{
          setTimeout(check, 200);
          return;
        }}
        var hasCharts = root.querySelectorAll('canvas').length > 0 ||
                        root.querySelectorAll('[_echarts_instance_]').length > 0;
        setTimeout(function() {{
          window.resizeAllCharts && window.resizeAllCharts();
          window.__ARTIFACT_RENDER_COMPLETE__ = true;
        }}, hasCharts ? 1500 : 300);
      }}
      setTimeout(check, 200);
    }})();
  {SC}
</body>
</html>"""
        return html

    # Maximum in-tool repair LLM calls per artifact operation. Bounded so a
    # stubborn defect can't consume the whole tool timeout — after this the
    # tool returns a structured failure and the planner decides.
    MAX_RENDER_REPAIR_ATTEMPTS = 5

    @staticmethod
    def fatal_render_errors(errors: List[str]) -> List[str]:
        """Errors that gate completion: thrown/uncaught exceptions.

        `[console.error]` entries are advisory — React dev builds log
        non-fatal warnings through console.error, so they inform repair but
        never fail an otherwise working artifact.
        """
        return [e for e in (errors or []) if not e.startswith("[console.error]")]

    @staticmethod
    def params_wiring_errors(code: str, artifact_data: Dict[str, Any]) -> List[str]:
        """Contract check: declared input params must be WIRED in the code.

        A dashboard whose queries declare input parameters but whose code never
        calls useParams()/setParam renders controls that only mutate local
        state — selecting a value silently never re-runs the data. Returns
        synthetic repair errors (empty when satisfied) that ride the same
        repair loop as render errors. Identity params are exempt: they render
        as a badge, not a control.
        """
        names: List[str] = []
        seen: set = set()
        for v in (artifact_data or {}).get("visualizations") or []:
            for p in (v.get("parameters") or []) if isinstance(v, dict) else []:
                name = p.get("name") if isinstance(p, dict) else None
                if not name or name in seen or (p.get("source") or "input") == "identity":
                    continue
                seen.add(name)
                names.append(name)
        if not names:
            return []
        src = code or ""
        if "useParams" not in src or "setParam" not in src:
            return [
                "[params contract] This dashboard's queries declare input parameter(s) "
                f"{', '.join(names)}, but the code never calls useParams()/setParam — a "
                "control that only sets local React state never re-runs the data. Wire "
                "every param control to useParams().setParam('<name>', value) — bind "
                "option.value (never the label), drive a loading state from "
                "useParams().loading, and render useParams().error when set."
            ]
        missing = [n for n in names if n not in src]
        if missing:
            return [
                "[params contract] Parameter(s) " + ", ".join(missing) + " are declared "
                "by this dashboard's queries but never referenced in the code — each "
                "needs a control wired to useParams().setParam('<name>', value)."
            ]
        return []

    async def _fix_code(
        self,
        code: str,
        errors: List[str],
        mode: str,
        runtime_ctx: Dict[str, Any],
    ) -> Optional[str]:
        """Compact in-tool repair: current code + exact errors → corrected code.

        Deliberately does NOT rebuild the full generation prompt — the model
        already produced this code; it needs the error, not the whole context.
        Returns the corrected code, or None if repair was unavailable/failed.
        """
        sigkill_event = runtime_ctx.get("sigkill_event")
        if sigkill_event and sigkill_event.is_set():
            return None

        error_text = "\n".join(f"- {e}" for e in errors[:5])

        fix_prompt = f"""You previously wrote the React dashboard code below. It runs in a sandboxed iframe (React 18 + Babel standalone + Tailwind + ECharts via <EChart>, data via useArtifactData()). When rendered, it produced these errors:

{error_text}

Current code:
```
{code}
```

Fix ONLY what the errors require — do not redesign, restyle, or restructure anything else. Common causes: using a variable/component before its declaration, duplicate declarations, undefined symbols, unguarded nullish values before string methods, malformed JSX.

Output the FULL corrected code wrapped in <script type="text/babel"> ... </script>. No explanations, no diff markers, no markdown fences."""

        llm = LLM(runtime_ctx.get("model"), usage_session_maker=async_session_maker)
        try:
            chunks: list[str] = []
            async for evt in llm.inference_stream_v2(
                messages=[Message(role="user", content=fix_prompt)],
                usage_scope="create_artifact_fix",
            ):
                if sigkill_event and sigkill_event.is_set():
                    return None
                if isinstance(evt, TextDeltaEvent):
                    chunks.append(evt.text)
            fixed = self._extract_code("".join(chunks), mode=mode)
            return fixed if fixed and fixed.strip() else None
        except Exception:
            logger.exception("In-tool code repair failed")
            return None

    async def _fix_pptx_code(
        self,
        code: str,
        error: str,
        runtime_ctx: Dict[str, Any],
    ) -> Optional[str]:
        """Compact in-tool repair for slides: current code + exact error → corrected code.

        Slides twin of _fix_code: the model already produced this code; it
        needs the executor's error, not the whole generation context.
        Returns the corrected code, or None if repair was unavailable/failed.
        """
        sigkill_event = runtime_ctx.get("sigkill_event")
        if sigkill_event and sigkill_event.is_set():
            return None

        fix_prompt = f"""You previously wrote the python-pptx code below. It runs in a sandboxed namespace that already provides: Presentation, Inches, Pt, Emu, RGBColor, PP_ALIGN, MSO_ANCHOR, MSO_SHAPE, XL_CHART_TYPE, XL_LEGEND_POSITION, CategoryChartData, ChartData, plus the data variables `visualizations` (list of dicts with 'title', 'columns', 'rows'), `report`, `image`/`image_ids`, and `_pptx_output_path` (the code must call prs.save(_pptx_output_path)). Note that each entry of viz['columns'] is a DICT like {{'field': 'Revenue', 'headerName': 'Revenue'}} — use col['field'] to get the row key, never pass the dict itself where a string is expected. When executed, it failed with this error:

{error}

Current code:
```python
{code}
```

Fix ONLY what the error requires — do not redesign, restyle, or restructure anything else. Common causes: passing a non-string (e.g. a column dict) to a text property, indexing a row with the wrong key, referencing an undefined name, wrong python-pptx API usage.

Output the FULL corrected code in a ```python code block. No explanations, no diff markers."""

        llm = LLM(runtime_ctx.get("model"), usage_session_maker=async_session_maker)
        try:
            chunks: list[str] = []
            async for evt in llm.inference_stream_v2(
                messages=[Message(role="user", content=fix_prompt)],
                usage_scope="create_artifact_fix",
            ):
                if sigkill_event and sigkill_event.is_set():
                    return None
                if isinstance(evt, TextDeltaEvent):
                    chunks.append(evt.text)
            fixed = self._extract_code("".join(chunks), mode="slides")
            return fixed if fixed and fixed.strip() else None
        except Exception:
            logger.exception("In-tool PPTX code repair failed")
            return None

    @staticmethod
    def _pptx_error_text(exc: Exception) -> str:
        """Trimmed executor error for repair prompts and observations.

        Keeps the frames inside the generated code (reported as <string>) and
        the final exception line — the sandbox internals above them are noise
        the repair model can't act on.
        """
        import traceback as _tb

        lines = _tb.format_exception(type(exc), exc, exc.__traceback__)
        flat = "".join(lines).splitlines()
        kept = [l for l in flat if '<string>' in l]
        # Final line always carries the exception type + message
        if flat:
            kept.append(flat[-1])
        text = "\n".join(kept[-12:]) if kept else str(exc)
        return text[:2000]

    @staticmethod
    def _first_preview_base64(preview_images: List[str]) -> Optional[str]:
        """Base64 of the first slide preview PNG, or None if unavailable."""
        if not preview_images:
            return None
        import base64

        path = Path(__file__).parent.parent.parent.parent.parent / "uploads" / preview_images[0]
        try:
            return base64.b64encode(path.read_bytes()).decode("ascii")
        except OSError:
            return None

    async def _execute_and_repair_pptx(
        self,
        code: str,
        visualizations: List[Dict[str, Any]],
        report_data: Dict[str, Any],
        output_path: Path,
        images: Dict[str, bytes],
        runtime_ctx: Dict[str, Any],
        deadline_monotonic: Optional[float] = None,
    ) -> AsyncIterator[Any]:
        """Execute slides code and repair it in-tool when it fails.

        Slides twin of _validate_and_repair_stream. Async generator: yields
        ToolProgressEvent items for the UI, then a final dict result:
          {"code", "ok", "error", "repair_attempts"}

        - Only executor exceptions gate; a successful run ends the loop.
        - If repair can't produce a working deck, the ORIGINAL code and its
          error are returned so what's persisted matches what was executed.
        """
        import time as _time

        def _time_left() -> bool:
            return deadline_monotonic is None or _time.monotonic() < deadline_monotonic

        sigkill_event = runtime_ctx.get("sigkill_event")
        executor = PptxCodeExecutor(logger=logger)

        def _try(candidate_code: str) -> Optional[str]:
            """Run once; returns None on success, trimmed error text on failure."""
            try:
                executor.execute_pptx_code(
                    code=candidate_code,
                    visualizations=visualizations,
                    report=report_data,
                    output_path=output_path,
                    images=images,
                )
                return None
            except Exception as e:
                logger.error(f"PPTX execution failed: {e}")
                return self._pptx_error_text(e)

        yield ToolProgressEvent(type="tool.progress", payload={"stage": "executing_pptx_code"})
        error = _try(code)

        original_code = code
        original_error = error
        candidate = code
        attempts = 0

        while (
            error
            and attempts < self.MAX_RENDER_REPAIR_ATTEMPTS
            and _time_left()
            and not (sigkill_event and sigkill_event.is_set())
        ):
            attempts += 1
            yield ToolProgressEvent(
                type="tool.progress",
                payload={"stage": "repairing_code", "attempt": attempts, "error_count": 1},
            )
            fixed = await self._fix_pptx_code(candidate, error, runtime_ctx)
            if not fixed or fixed == candidate:
                break

            yield ToolProgressEvent(
                type="tool.progress",
                payload={"stage": "executing_pptx_code", "attempt": attempts},
            )
            error = _try(fixed)
            candidate = fixed

        if not error:
            yield {
                "code": candidate,
                "ok": True,
                "error": None,
                "repair_attempts": attempts,
            }
        else:
            # Repair didn't converge — return the original, verified-broken
            # state rather than an unverified intermediate, and drop any deck
            # a partial attempt may have left behind.
            try:
                if output_path.exists():
                    output_path.unlink()
            except OSError:
                pass
            yield {
                "code": original_code,
                "ok": False,
                "error": original_error,
                "repair_attempts": attempts,
            }

    async def _validate_and_repair_stream(
        self,
        code: str,
        artifact_data: Dict[str, Any],
        mode: str,
        runtime_ctx: Dict[str, Any],
        deadline_monotonic: Optional[float] = None,
    ) -> AsyncIterator[Any]:
        """Render-validate page code and repair it in-tool when it fails.

        Async generator: yields ToolProgressEvent items for the UI, then a
        final dict result:
          {"code", "clean", "screenshot", "errors", "repair_attempts"}

        - Validation always runs (it needs Playwright, not a vision model or
          the allow_llm_see_data flag — errors are not data).
        - Only fatal (thrown) errors gate; console errors ride along as
          repair context.
        - If repair can't produce a clean render, the ORIGINAL code and its
          errors are returned so what's persisted matches what was verified.
        """
        import time as _time

        def _time_left() -> bool:
            return deadline_monotonic is None or _time.monotonic() < deadline_monotonic

        sigkill_event = runtime_ctx.get("sigkill_event")

        yield ToolProgressEvent(type="tool.progress", payload={"stage": "validating_render"})
        try:
            html = self._build_thumbnail_html(artifact_data, code, mode=mode)
        except Exception as e:
            # Validation infrastructure unavailable (e.g. vendored libs not
            # downloaded) must not fail artifact creation — pass through
            # unvalidated, exactly like the missing-Playwright fallback.
            logger.warning(f"Render validation unavailable (thumbnail HTML build failed): {e}")
            yield {
                "code": code,
                "clean": True,
                "screenshot": None,
                "errors": [],
                "repair_attempts": 0,
                "validation_skipped": True,
            }
            return
        screenshot, errors = await self._take_preview_screenshot(html)
        fatal = self.fatal_render_errors(errors)
        # Params-wiring contract rides the same repair loop: an unwired param
        # renders fine but silently never re-runs the data — as much a defect
        # as a thrown error, and fixable by the same in-tool repair.
        wiring = self.params_wiring_errors(code, artifact_data) if mode == "page" else []

        original_code = code
        original_screenshot = screenshot
        original_errors = errors
        original_wiring = list(wiring)
        candidate = code
        attempts = 0

        while (
            (fatal or wiring)
            and attempts < self.MAX_RENDER_REPAIR_ATTEMPTS
            and _time_left()
            and not (sigkill_event and sigkill_event.is_set())
        ):
            attempts += 1
            yield ToolProgressEvent(
                type="tool.progress",
                payload={"stage": "repairing_code", "attempt": attempts, "error_count": len(fatal) + len(wiring)},
            )
            fixed = await self._fix_code(candidate, list(errors) + wiring, mode, runtime_ctx)
            if not fixed or fixed == candidate:
                break

            yield ToolProgressEvent(
                type="tool.progress",
                payload={"stage": "validating_render", "attempt": attempts},
            )
            html = self._build_thumbnail_html(artifact_data, fixed, mode=mode)
            screenshot, errors = await self._take_preview_screenshot(html)
            candidate = fixed
            fatal = self.fatal_render_errors(errors)
            wiring = self.params_wiring_errors(candidate, artifact_data) if mode == "page" else []

        if not fatal:
            # Render fatality decides candidate vs original; unconverged wiring
            # only ANNOTATES — a rendering dashboard with a dead control still
            # beats no dashboard, but the caller must be able to say so.
            yield {
                "code": candidate,
                "clean": True,
                "screenshot": screenshot,
                "errors": errors,
                "repair_attempts": attempts,
                "params_wiring_errors": wiring,
            }
        else:
            # Repair didn't converge — return the original, verified-broken
            # state rather than an unverified intermediate.
            yield {
                "code": original_code,
                "clean": False,
                "screenshot": original_screenshot,
                "errors": original_errors,
                "repair_attempts": attempts,
                "params_wiring_errors": original_wiring,
            }

    def _build_viz_profile(self, viz: Dict[str, Any], allow_llm_see_data: bool) -> Dict[str, Any]:
        """Build a privacy-aware profile of a visualization's data."""
        # Enrich columns with dtype/unique_count/min/max from column_info (always — not sensitive)
        column_info = viz.get("column_info") or {}
        raw_columns = viz.get("columns", [])
        enriched_columns = []
        for c in raw_columns:
            col = dict(c) if isinstance(c, dict) else {"field": c}
            field = col.get("field") or col.get("headerName") or col.get("name")
            if field and field in column_info:
                meta = column_info[field]
                col["dtype"] = meta.get("dtype")
                col["unique_count"] = meta.get("unique_count")
                if meta.get("min") is not None:
                    col["min"] = meta["min"]
                if meta.get("max") is not None:
                    col["max"] = meta["max"]
            enriched_columns.append(col)

        profile: Dict[str, Any] = {
            "id": viz.get("id"),
            "title": viz.get("title"),
            "query_id": viz.get("query_id"),
            "chart_type": viz.get("data_model_type") or "table",
            # True dataset size; sample_row_count is how many rows are shown
            # to generation/preview (capped). At runtime the dashboard
            # receives ALL row_count rows.
            "row_count": viz.get("row_count", 0),
            "sample_row_count": viz.get("sample_row_count", len(viz.get("rows") or [])),
            "columns": enriched_columns,
        }

        # Include data model hints
        data_model = viz.get("dataModel") or {}
        if data_model:
            series = data_model.get("series", [])
            if series:
                profile["series_config"] = series[:3]  # First 3 series configs
            if data_model.get("group_by"):
                profile["group_by"] = data_model.get("group_by")

        # Include view configuration hints
        view = viz.get("view") or {}
        if view:
            inner_view = view.get("view") or view
            profile["view_config"] = {
                "type": inner_view.get("type"),
                "x": inner_view.get("x"),
                "y": inner_view.get("y"),
                "category": inner_view.get("category"),
                "value": inner_view.get("value"),
            }
            # Surface aggregation (top-level) + per-series aggregations so the
            # artifact can honor granular-data handling rather than reading
            # the first row.
            if inner_view.get("aggregation"):
                profile["view_config"]["aggregation"] = inner_view.get("aggregation")
            series_styles = inner_view.get("seriesStyles") or []
            series_aggs = [
                {"key": s.get("key"), "aggregation": s.get("aggregation")}
                for s in series_styles
                if isinstance(s, dict) and s.get("aggregation")
            ]
            if series_aggs:
                profile["view_config"]["series_aggregations"] = series_aggs
            default_filters = inner_view.get("defaultFilters") or []
            if default_filters:
                profile["view_config"]["default_filters"] = default_filters
            # Include palette if present
            palette = inner_view.get("palette") or {}
            if palette.get("colors"):
                profile["colors"] = palette.get("colors")[:5]

        # Declared server-side query parameters (useParams contract). Compact:
        # the artifact needs name/type/source/label/default/options to build
        # controls, nothing more.
        params = viz.get("parameters") or []
        if params:
            profile["parameters"] = [
                {
                    k: p.get(k)
                    for k in ("name", "type", "label", "source", "default", "required", "options")
                    if p.get(k) is not None or k in ("name", "source")
                }
                for p in params
                if isinstance(p, dict) and p.get("name")
            ]

        # Include sample data if allowed
        if allow_llm_see_data:
            rows = viz.get("rows", [])
            if rows:
                profile["sample_rows"] = rows[:5]  # First 5 rows
                # Compute basic stats for numeric columns
                if rows and isinstance(rows[0], dict):
                    stats = {}
                    for col in viz.get("columns", []):
                        col_name = col if isinstance(col, str) else col.get("field", col.get("name"))
                        if col_name:
                            values = [r.get(col_name) for r in rows if r.get(col_name) is not None]
                            numeric_values = [v for v in values if isinstance(v, (int, float))]
                            if numeric_values:
                                stats[col_name] = {
                                    "min": min(numeric_values),
                                    "max": max(numeric_values),
                                    "sample_values": numeric_values[:3]
                                }
                            elif values:
                                unique = list(set(str(v) for v in values[:20]))
                                stats[col_name] = {
                                    "unique_count": len(unique),
                                    "sample_values": unique[:5]
                                }
                    if stats:
                        profile["column_stats"] = stats

        return profile

    async def run_stream(self, tool_input: Dict[str, Any], runtime_ctx: Dict[str, Any]) -> AsyncIterator[ToolEvent]:
        data = CreateArtifactInput(**tool_input)
        # Repair budget: leave headroom under the runner's 300s hard timeout
        # for the final persist + observation after the last repair round.
        _repair_deadline = time.monotonic() + 210

        # Early validation: require at least one visualization OR at least one file
        # (an image/PDF-only artifact is allowed when file_ids are provided).
        # Slides are exempt: a deck legitimately opens with a title, agenda or
        # narrative slide that carries no chart, and a whole deck may be
        # narrative-only.
        if (
            (not data.visualization_ids or len(data.visualization_ids) == 0)
            and not getattr(data, "file_ids", None)
            and data.mode != "slides"
        ):
            yield ToolStartEvent(type="tool.start", payload={"title": data.title or "Artifact"})
            yield ToolEndEvent(
                type="tool.end",
                payload={
                    "output": {
                        "success": False,
                        "error": "No visualization_ids provided. At least one visualization is required to create an artifact.",
                    },
                    "observation": {
                        "summary": "Failed to create artifact: no visualization_ids provided",
                        "error": {
                            "type": "validation_error",
                            "message": "visualization_ids is required and must contain at least one visualization ID. Create visualizations using create_data first, then use their IDs here.",
                        },
                    },
                },
            )
            return

        yield ToolStartEvent(type="tool.start", payload={"title": data.title or "Artifact"})
        yield ToolProgressEvent(type="tool.progress", payload={"stage": "init"})

        # Get runtime context
        sigkill_event = runtime_ctx.get("sigkill_event")
        report = runtime_ctx.get("report")
        user = runtime_ctx.get("user")
        organization = runtime_ctx.get("organization")
        db = runtime_ctx.get("db")
        context_hub = runtime_ctx.get("context_hub")
        organization_settings = runtime_ctx.get("settings")

        # Check privacy setting
        allow_llm_see_data = True
        if organization_settings:
            try:
                allow_llm_see_data = organization_settings.get_config("allow_llm_see_data").value
            except Exception:
                allow_llm_see_data = True

        instruction_context_builder = runtime_ctx.get("instruction_context_builder") or (
            getattr(context_hub, "instruction_builder", None) if context_hub else None
        )

        # Get conversation history context (similar to create_data.py)
        context_view = runtime_ctx.get("context_view")
        messages_context = ""
        try:
            _messages_section_obj = getattr(context_view.warm, "messages", None) if context_view else None
            messages_context = _messages_section_obj.render() if _messages_section_obj else ""
        except Exception as e:
            logger.warning(f"Failed to extract messages context: {e}")
            messages_context = ""

        # Load images attached to the head completion for vision-capable models
        head_completion = runtime_ctx.get("head_completion")
        head_completion_id = str(head_completion.id) if head_completion else None
        completion_images = await self._load_completion_images(db, head_completion_id)

        # Validate model supports vision if images are present
        model = runtime_ctx.get("model")
        if completion_images and not getattr(model, "supports_vision", False):
            logger.info(f"Model doesn't support vision, skipping {len(completion_images)} completion images")
            completion_images = []

        # Note: Previous artifacts are now available via observation context (from create_artifact/read_artifact)
        # No need to fetch from DB - the planner can call read_artifact if needed

        # Fetch visualizations by ID from database
        visualizations: List[Dict[str, Any]] = []
        warnings: List[str] = []
        included_viz_ids: List[str] = []

        # Fetch all visualizations in a single batched query
        yield ToolProgressEvent(type="tool.progress", payload={"stage": "loading_visualizations"})
        from app.models.query import Query
        from app.models.step import Step
        report_id = str(report.id) if report else None
        try:
            # populate_existing=True forces SQLAlchemy to refresh objects from DB
            # rather than returning stale identity-map copies (e.g. query.steps or
            # query.default_step may have been loaded before the step was created/updated)
            result = await db.execute(
                select(Visualization)
                .options(
                    selectinload(Visualization.query).selectinload(Query.default_step),
                    selectinload(Visualization.query).selectinload(Query.steps),
                )
                .where(Visualization.id.in_(data.visualization_ids))
                .execution_options(populate_existing=True)
            )
            fetched_vizs = {str(v.id): v for v in result.scalars().all()}
        except Exception as e:
            logger.exception("Failed to batch-fetch visualizations")
            fetched_vizs = {}
            warnings.append(f"Error fetching visualizations: {str(e)}")

        # Process each requested viz in order, validating and building entries
        for viz_id in data.visualization_ids:
            viz = fetched_vizs.get(viz_id)
            if viz is None:
                warnings.append(f"Visualization {viz_id} not found")
                continue

            # Validate viz belongs to the report
            if report_id and str(viz.report_id) != report_id:
                warnings.append(f"Visualization {viz_id} does not belong to this report")
                continue

            # Get the step with data (prefer default_step, fallback to latest step)
            step = None
            if viz.query and viz.query.default_step:
                step = viz.query.default_step
            elif viz.query and viz.query.steps:
                step = viz.query.steps[-1] if viz.query.steps else None

            # Check if the associated step is successful
            step_status = step.status if step else None
            if step_status != "success":
                _has_query = viz.query is not None
                _has_default = viz.query.default_step is not None if _has_query else False
                _steps_len = len(viz.query.steps) if _has_query and viz.query.steps else 0
                _default_step_id = getattr(viz.query, 'default_step_id', None) if _has_query else None
                logger.warning(
                    f"Visualization {viz_id} skipped: step_status='{step_status}', "
                    f"has_query={_has_query}, default_step_id={_default_step_id}, "
                    f"has_default_step={_has_default}, steps_count={_steps_len}"
                )
                warnings.append(f"Visualization {viz_id} skipped: step status is '{step_status or 'unknown'}' (not success)")
                continue

            # Get data directly from step (like frontend does). Generation and
            # headless preview see a bounded sample; the live render receives
            # the full dataset from the frontend. Keep both counts so the
            # sample is never mislabeled as the whole dataset.
            step_data = step.data if step else {}
            _all_rows = (step_data.get("rows") or []) if step_data else []
            rows = _all_rows[:100]
            total_row_count = len(_all_rows)
            raw_columns = step_data.get("columns") or [] if step_data else []
            data_model = step.data_model if step else {}
            step_info = step_data.get("info") or {} if step_data else {}
            column_info = step_info.get("column_info") or {}

            # Keep raw column objects (with field/headerName) — matches the prompt contract
            columns = raw_columns

            # Extract field names for internal use (filterable columns, logging)
            column_fields = []
            for c in raw_columns:
                if isinstance(c, str):
                    column_fields.append(c)
                elif isinstance(c, dict):
                    col_name = c.get("field") or c.get("colId") or c.get("headerName") or c.get("name")
                    if col_name:
                        column_fields.append(col_name)

            # Build visualization entry
            view_dict = viz.view or {}
            query_id = str(viz.query_id) if viz.query_id else None

            ventry = {
                "id": str(viz.id),
                "title": viz.title,
                "query_id": query_id,
                "view": self._trim_none(view_dict),
                "data_model_type": (view_dict.get("view") or {}).get("type") or view_dict.get("type"),
                "columns": columns,
                "column_info": column_info,
                # row_count is the TRUE dataset size; rows is a sample capped
                # at 100 for generation/preview. The old code reported the
                # truncated length as row_count, so the LLM couldn't know the
                # data was sampled and generated truncation workarounds.
                "row_count": total_row_count,
                "sample_row_count": len(rows),
                "rows": rows,
                # The profile the planner reads calls this same list
                # `sample_rows` (see _build_viz_profile), so generated code
                # reaches for either name. Expose both: picking the wrong one
                # used to yield [] and fail the whole deck with
                # "chart data contains no categories".
                "sample_rows": rows,
                "dataModel": data_model or {},
                # Declared query parameters (ParamSpec dicts): the dashboard
                # should render a control per input param via useParams() and
                # a "scoped to you" badge for identity params.
                "parameters": list(getattr(viz.query, "parameters", None) or []) if viz.query else [],
            }

            # Debug logging
            logger.info(f"Visualization {viz.title}: {len(rows)} rows, {len(column_fields)} columns: {column_fields[:5] if column_fields else 'none'}")
            if rows:
                logger.info(f"  Sample row keys: {list(rows[0].keys())[:5] if isinstance(rows[0], dict) else 'not a dict'}")

            visualizations.append(ventry)
            included_viz_ids.append(str(viz.id))

        # Resolve any embedded files (generated images / uploaded images or PDFs).
        # Scoped to the org; stored on the artifact content so the frontend can
        # fetch + inject them into the sandbox for the <BowFile> component.
        included_files: List[Dict[str, Any]] = []
        requested_file_ids = getattr(data, "file_ids", None) or []
        if requested_file_ids:
            try:
                file_result = await db.execute(
                    select(File).where(
                        File.id.in_([str(f) for f in requested_file_ids]),
                        File.organization_id == str(organization.id) if organization else File.organization_id.is_(None),
                    )
                )
                fetched_files = {str(f.id): f for f in file_result.scalars().all()}
            except Exception as e:
                logger.warning(f"create_artifact: failed to fetch files: {e}")
                fetched_files = {}
            for fid in requested_file_ids:
                f = fetched_files.get(str(fid))
                if f is None:
                    warnings.append(f"File {fid} not found or not in this organization")
                    continue
                included_files.append({
                    "id": str(f.id),
                    "content_type": f.content_type or "application/octet-stream",
                    "filename": f.filename,
                })

        # Early failure: if no valid visualizations AND no files were resolved,
        # fail like create_data does with tables. Slides may be narrative-only
        # (see the mode exemption in the early validation above).
        if not visualizations and not included_files and data.mode != "slides":
            yield ToolEndEvent(
                type="tool.end",
                payload={
                    "output": {
                        "success": False,
                        "error": "No valid visualizations found. All requested visualization_ids were either not found, don't belong to this report, or have non-success step status.",
                    },
                    "observation": {
                        "summary": "Failed to create artifact: no valid visualizations resolved",
                        "error": {
                            "type": "no_valid_visualizations",
                            "message": "None of the requested visualization_ids could be used. Ensure visualizations exist, belong to this report, and have successful step status.",
                            "requested_ids": data.visualization_ids,
                            "warnings": warnings,
                        },
                    },
                },
            )
            return

        # Build visualization profiles (privacy-aware)
        yield ToolProgressEvent(type="tool.progress", payload={"stage": "building_profiles"})
        viz_profiles = [self._build_viz_profile(v, allow_llm_see_data) for v in visualizations]

        # Emit visualizations_resolved
        yield ToolProgressEvent(type="tool.progress", payload={
            "stage": "visualizations_resolved",
            "tool_name": "create_artifact",
            "visualizations": [
                {"id": v["id"], "title": v["title"], "type": v.get("data_model_type", "")}
                for v in visualizations
            ],
        })

        # Build instruction context
        yield ToolProgressEvent(type="tool.progress", payload={"stage": "building_context"})
        instructions_context = ""
        try:
            if instruction_context_builder is not None:
                inst_section = await instruction_context_builder.build(categories=["dashboard", "visualization", "general"])
                instructions_context = inst_section.render() or ""
        except Exception:
            pass

        # Create artifact early with pending status so frontend can show it
        artifact = Artifact(
            report_id=str(report.id) if report else None,
            user_id=str(user.id) if user else None,
            organization_id=str(organization.id) if organization else None,
            title=data.title or "Untitled Artifact",
            mode=data.mode,
            content={},  # Empty content initially
            generation_prompt=data.prompt,
            version=1,
            status="pending",
        )
        db.add(artifact)
        await db.commit()
        await db.refresh(artifact)

        # Notify frontend that artifact is created (pending)
        yield ToolProgressEvent(
            type="tool.progress",
            payload={
                "stage": "artifact_created",
                "artifact_id": str(artifact.id),
                "status": "pending",
                "timing": False,
            }
        )

        # Build the prompt for generating React code
        yield ToolProgressEvent(type="tool.progress", payload={"stage": "building_prompt"})

        # Identity vocabulary for page mode: the requester's current_user as an
        # example + org group names, so identity intent ("show to Finance",
        # "greet by department") resolves against real names instead of guesses.
        identity_context = ""
        if data.mode == "page":
            identity_context = await build_identity_context(db, user, organization)

        prompt = self._build_prompt(
            user_prompt=data.prompt,
            title=data.title,
            mode=data.mode,
            viz_profiles=viz_profiles,
            instructions_context=instructions_context,
            identity_context=identity_context,
            report_title=getattr(report, 'title', None) if report else None,
            allow_llm_see_data=allow_llm_see_data,
            messages_context=messages_context,
            image_count=len(completion_images),
            organization_settings=organization_settings,
            files=included_files,
        )
        # Static reference goes in the system prompt so provider-side prompt
        # caching reuses it across artifact calls (page mode only — slides
        # keeps its single-prompt path).
        system_prompt = self._build_page_system_prompt() if data.mode == "page" else None

        # Stream from LLM
        yield ToolProgressEvent(type="tool.progress", payload={"stage": "llm_generating"})
        llm = LLM(runtime_ctx.get("model"), usage_session_maker=async_session_maker)
        buffer = ""
        slides_detected = 0  # Track number of slides detected during streaming

        async for evt in llm.inference_stream_v2(
            messages=[Message(role="user", content=prompt)],
            system=system_prompt,
            images=completion_images if completion_images else None,
            usage_scope="create_artifact",
            usage_scope_ref_id=str(report.id) if report else None,
        ):
            if sigkill_event and sigkill_event.is_set():
                break
            if isinstance(evt, TextDeltaEvent):
                buffer += evt.text

            # For slides mode, detect new slides as they're generated
            if data.mode == "slides":
                # Count slide sections in buffer
                current_slides = buffer.count('<section class="slide"')
                if current_slides > slides_detected:
                    # New slide detected
                    for i in range(slides_detected, current_slides):
                        yield ToolProgressEvent(
                            type="tool.progress",
                            payload={
                                "stage": "slide_generated",
                                "slide_index": i,
                                "total_slides": current_slides,
                                "timing": False,
                            }
                        )
                    slides_detected = current_slides

            # Stream partial updates
            if len(buffer) % 100 == 0:  # Throttle updates
                yield ToolProgressEvent(
                    type="tool.progress",
                    payload={"stage": "generating", "chars": len(buffer), "timing": False}
                )

        # Check sigkill after LLM generation
        if sigkill_event and sigkill_event.is_set():
            # Update artifact to stopped status
            artifact.status = "stopped"
            await db.commit()
            yield ToolEndEvent(
                type="tool.end",
                payload={
                    "output": {"success": False, "artifact_id": str(artifact.id), "error": "Stopped by user"},
                    "observation": {"summary": "Artifact creation stopped by user", "artifact_id": str(artifact.id), "stopped": True},
                },
            )
            return

        # Extract the code from the response
        code = self._extract_code(buffer, mode=data.mode)

        # ═══════════════════════════════════════════════════════════════════════
        # Mode-specific processing: slides uses python-pptx, page skips to save
        # ═══════════════════════════════════════════════════════════════════════

        pptx_path: Optional[str] = None
        pptx_success: bool = True
        pptx_error: Optional[str] = None
        pptx_repair_attempts: int = 0
        preview_images: List[str] = []

        if data.mode == "slides":
            # ═══════════════════════════════════════════════════════════════════
            # SLIDES MODE: Execute python-pptx code (repairing in-tool on
            # failure, mirroring page-mode render validation) and generate
            # previews.
            # ═══════════════════════════════════════════════════════════════════
            report_data = {
                "id": str(report.id) if report else None,
                "title": getattr(report, "title", None) if report else None,
                "theme": getattr(report, "theme", None) if report else None,
            }

            uploads_dir = Path(__file__).parent.parent.parent.parent.parent / "uploads" / "pptx"
            uploads_dir.mkdir(parents=True, exist_ok=True)
            output_path = uploads_dir / f"{artifact.id}.pptx"

            _pptx_result: Optional[Dict[str, Any]] = None
            async for _item in self._execute_and_repair_pptx(
                code,
                visualizations,
                report_data,
                output_path,
                await load_image_bytes(db, included_files),
                runtime_ctx,
                deadline_monotonic=_repair_deadline,
            ):
                if isinstance(_item, dict):
                    _pptx_result = _item
                else:
                    yield _item
            if _pptx_result is not None:
                code = _pptx_result["code"]
                pptx_success = bool(_pptx_result["ok"])
                pptx_error = _pptx_result["error"]
                pptx_repair_attempts = int(_pptx_result["repair_attempts"] or 0)
            if pptx_success:
                pptx_path = str(output_path)

            # Previews are rendered by LibreOffice, which is a separate concern
            # from building the deck: it can be missing an import filter or be
            # misconfigured while the .pptx itself is perfectly valid. Failing
            # the artifact here would also make the export endpoint refuse to
            # serve a deck the user can open, so a preview failure only costs
            # the preview.
            if pptx_success and pptx_path:
                yield ToolProgressEvent(
                    type="tool.progress",
                    payload={"stage": "generating_previews"}
                )
                try:
                    preview_service = PptxPreviewService(logger=logger)
                    preview_images = preview_service.generate_previews(
                        pptx_path=Path(pptx_path),
                        artifact_id=str(artifact.id),
                    )
                except Exception as e:
                    logger.warning(
                        f"PPTX preview generation failed; deck is still downloadable: {e}"
                    )

        # ═══════════════════════════════════════════════════════════════════════
        # Page mode: render-validate (and repair in-tool) BEFORE persisting.
        # "completed" must mean "renders without fatal errors" — the old flow
        # committed completed first and only then discovered render errors,
        # pushing every repair through a full outer planner iteration.
        # Validation runs regardless of vision support or allow_llm_see_data:
        # errors are not data. Only the screenshot ATTACHMENT is gated below.
        # ═══════════════════════════════════════════════════════════════════════
        screenshot_base64: Optional[str] = None
        render_errors: list[str] = []
        render_clean = True
        repair_attempts = 0
        params_wiring_errors: list[str] = []
        thumbnail_html: Optional[str] = None
        artifact_data: Optional[Dict[str, Any]] = None

        if data.mode == "page":
            artifact_data = {
                "report": {
                    "id": str(report.id) if report else None,
                    "title": getattr(report, "title", None) if report else None,
                    "theme": getattr(report, "theme", None) if report else None,
                },
                "visualizations": visualizations,
                # Headless validation renders anonymously on purpose: identity
                # is per-viewer, injected by the host at render time, and this
                # exercises the null-guard path in every artifact before it is
                # persisted.
                "current_user": None,
            }
            # Inline embedded files as data URIs so the headless render
            # (which has no auth context) can show images/PDFs via <BowFile>.
            if included_files:
                artifact_data["files"] = await self._build_file_datauris(db, included_files)

            _validate_result: Optional[Dict[str, Any]] = None
            async for _item in self._validate_and_repair_stream(
                code, artifact_data, data.mode, runtime_ctx, deadline_monotonic=_repair_deadline,
            ):
                if isinstance(_item, dict):
                    _validate_result = _item
                else:
                    yield _item
            if _validate_result is not None:
                code = _validate_result["code"]
                render_clean = bool(_validate_result["clean"])
                screenshot_base64 = _validate_result["screenshot"]
                render_errors = list(_validate_result["errors"] or [])
                repair_attempts = int(_validate_result["repair_attempts"] or 0)
                params_wiring_errors = list(_validate_result.get("params_wiring_errors") or [])
            try:
                thumbnail_html = self._build_thumbnail_html(artifact_data, code, mode=data.mode)
            except Exception as e:
                logger.warning(f"Thumbnail HTML build failed: {e}")
                thumbnail_html = None

        yield ToolProgressEvent(type="tool.progress", payload={"stage": "saving_artifact"})

        # Build content object (code is final — post-repair when repair ran)
        content: Dict[str, Any] = {
            "code": code,
            "visualization_ids": included_viz_ids,
        }

        # Embedded files (generated images / uploaded images/PDFs) referenced by
        # <BowFile id=...>. Stored as {id, content_type, filename}; the frontend
        # resolves the bytes and injects them into ARTIFACT_DATA.files.
        if included_files:
            content["files"] = included_files

        # Add slides-specific content
        if data.mode == "slides" and preview_images:
            content["preview_images"] = preview_images

        artifact.content = content
        if data.mode == "slides":
            artifact.status = "completed" if pptx_success else "failed"
            # Persist the executor error like page mode persists render errors,
            # so read_artifact and later repairs can see WHY the deck failed.
            if pptx_error:
                artifact.render_errors = [pptx_error]
        else:
            artifact.status = "completed" if render_clean else "failed"

        # Set pptx_path for slides mode
        if pptx_path:
            artifact.pptx_path = pptx_path

        # Persist screenshot and render errors for later retrieval (read_artifact)
        if screenshot_base64 or render_errors:
            artifact.screenshot_base64 = screenshot_base64
            artifact.render_errors = render_errors or None

        await db.commit()
        await db.refresh(artifact)

        if data.mode == "page" and thumbnail_html is not None and render_clean:
            # Generate thumbnail in background (for stored thumbnail, non-blocking)
            asyncio.create_task(
                self._generate_thumbnail_background(
                    artifact_id=str(artifact.id),
                    html_content=thumbnail_html,
                    mode=data.mode,
                )
            )
        elif preview_images:
            # For slides mode, use the first preview image as thumbnail
            first_preview = Path(__file__).parent.parent.parent.parent.parent / "uploads" / preview_images[0]
            if first_preview.exists():
                artifact.thumbnail_path = preview_images[0]
                await db.commit()

        # Slides mode that failed pptx execution (after bounded in-tool
        # repair): return a structured failure so the planner sees the real
        # error — mirror of the page-mode branch below. The artifact row is
        # persisted as status="failed" for debugging, but it is not presented
        # as a working deck.
        if data.mode == "slides" and not pptx_success:
            _error_msg = pptx_error or "unknown pptx execution error"
            slides_failure_observation: Dict[str, Any] = {
                "summary": (
                    f"Artifact '{data.title or 'Untitled'}' failed to build the presentation "
                    f"after {pptx_repair_attempts} in-tool repair attempt(s). "
                    f"Error: {_error_msg}"
                ),
                "error": {
                    "type": "pptx_execution_failed",
                    "message": _error_msg,
                    "repair_attempts": pptx_repair_attempts,
                    "remediation": (
                        "The generated python-pptx code does not execute. Call edit_artifact with an "
                        "edit_prompt that quotes the exact error above, or create_artifact to rebuild "
                        "with a simpler deck if the error persists."
                    ),
                },
                "artifact_id": str(artifact.id),
                "mode": data.mode,
                "visualization_count": len(visualizations),
                "visualization_ids": included_viz_ids,
                "render_errors": [_error_msg],
            }
            if warnings:
                slides_failure_observation["warnings"] = warnings
            yield ToolEndEvent(
                type="tool.end",
                payload={
                    "output": {
                        "success": False,
                        "artifact_id": str(artifact.id),
                        "error": f"PPTX execution failed: {_error_msg}",
                        "code_preview": {
                            "language": "python",
                            "code": code,
                            "collapsed_default": True,
                        },
                    },
                    "observation": slides_failure_observation,
                },
            )
            return

        # Page mode that failed render validation (after bounded in-tool
        # repair): return a structured failure so the planner sees the real
        # errors. The artifact row is persisted as status="failed" for
        # debugging, but it is not presented as a working dashboard.
        if data.mode == "page" and not render_clean:
            fatal_errors = self.fatal_render_errors(render_errors)
            first_error = fatal_errors[0] if fatal_errors else (render_errors[0] if render_errors else "unknown render error")
            failure_observation: Dict[str, Any] = {
                "summary": (
                    f"Artifact '{data.title or 'Untitled'}' failed render validation with "
                    f"{len(fatal_errors)} fatal error(s) after {repair_attempts} in-tool repair attempt(s). "
                    f"First error: {first_error}"
                ),
                "error": {
                    "type": "render_validation_failed",
                    "message": first_error,
                    "render_errors": render_errors,
                    "repair_attempts": repair_attempts,
                    "remediation": (
                        "The generated code does not render. Call edit_artifact with an edit_prompt "
                        "that quotes the exact error(s) above, or create_artifact to rebuild with a "
                        "simpler layout if the errors persist."
                    ),
                },
                "artifact_id": str(artifact.id),
                "mode": data.mode,
                "visualization_count": len(visualizations),
                "visualization_ids": included_viz_ids,
                "render_errors": render_errors,
            }
            if warnings:
                failure_observation["warnings"] = warnings
            _model = runtime_ctx.get("model")
            if screenshot_base64 and allow_llm_see_data and _model and getattr(_model, "supports_vision", False):
                failure_observation["images"] = [{
                    "data": screenshot_base64,
                    "media_type": "image/png",
                    "source_type": "base64",
                }]
            yield ToolEndEvent(
                type="tool.end",
                payload={
                    "output": {
                        "success": False,
                        "artifact_id": str(artifact.id),
                        "error": f"Render validation failed: {first_error}",
                        "code_preview": {
                            "language": "jsx",
                            "code": code,
                            "collapsed_default": True,
                        },
                    },
                    "observation": failure_observation,
                },
            )
            return

        output = CreateArtifactOutput(
            artifact_id=str(artifact.id),
            code=code,
            mode=data.mode,
            title=data.title,
            version=artifact.version,
        ).model_dump()

        # Add UI preview fields (similar to read_artifact)
        code_lines = code.count('\n') + 1 if code else 0
        output["artifact_preview"] = {
            "artifact_id": str(artifact.id),
            "title": data.title or "Untitled",
            "mode": data.mode,
            "version": artifact.version,
            "code_stats": {
                "chars": len(code),
                "lines": code_lines,
            },
            "visualization_ids": included_viz_ids,
            "visualization_count": len(visualizations),
        }
        # Code for collapsible toggle (collapsed by default in UI)
        output["code_preview"] = {
            "language": "jsx",
            "code": code,
            "collapsed_default": True,
        }

        # Build observation message
        summary_msg = f"Created artifact '{data.title or 'Untitled'}' with {len(code)} characters of code"
        if data.mode == "slides":
            if pptx_repair_attempts:
                summary_msg += f". PPTX execution passed after {pptx_repair_attempts} in-tool repair attempt(s)."
            if preview_images:
                summary_msg += f" Generated {len(preview_images)} slide preview images."
            else:
                summary_msg += (
                    " The deck was built, but slide preview images could not be generated on this "
                    "server — the PPTX file is still downloadable."
                )
        elif data.mode == "page":
            if repair_attempts:
                summary_msg += f". Render validation passed after {repair_attempts} in-tool repair attempt(s)."
            else:
                summary_msg += ". Render validation passed."
            console_warnings = [e for e in render_errors if e.startswith("[console.error]")]
            if console_warnings:
                summary_msg += f" {len(console_warnings)} non-fatal console error(s) were logged."

        # Screenshot attachment is gated on privacy + vision (the screenshot
        # shows the data); validation itself already ran unconditionally.
        _model = runtime_ctx.get("model")
        _attach_screenshot = bool(
            screenshot_base64 and allow_llm_see_data and _model and getattr(_model, "supports_vision", False)
        )
        if _attach_screenshot:
            summary_msg += " Screenshot of the rendered dashboard is attached — review it for visual correctness."
            summary_msg += ANON_PREVIEW_NOTE

        observation: Dict[str, Any] = {
            "summary": summary_msg,
            "artifact_id": str(artifact.id),
            "mode": data.mode,
            "visualization_count": len(visualizations),
            "visualization_ids": included_viz_ids,
        }
        if render_errors:
            observation["render_errors"] = render_errors
        if repair_attempts:
            observation["repair_attempts"] = repair_attempts
        # Wiring contract unmet after the repair budget: the dashboard renders
        # but its param control(s) are dead — say so instead of shipping it
        # silently, so the planner can follow up with edit_artifact.
        if params_wiring_errors:
            observation["params_wired"] = False
            observation["params_wiring_errors"] = params_wiring_errors
            observation["summary"] = (
                summary_msg
                + " WARNING: declared query parameters are NOT wired to controls — "
                "selecting a value will not re-run the data. Fix with edit_artifact: "
                + params_wiring_errors[0]
            )

        # Add preview screenshot for planner reflection (page mode)
        if _attach_screenshot:
            observation["images"] = [{
                "data": screenshot_base64,
                "media_type": "image/png",
                "source_type": "base64",
            }]

        # Add slides-specific info
        if data.mode == "slides":
            if preview_images:
                observation["preview_images"] = preview_images
                observation["slide_count"] = len(preview_images)
            if pptx_path:
                observation["pptx_path"] = pptx_path
            if pptx_repair_attempts:
                observation["repair_attempts"] = pptx_repair_attempts
            # Attach the first slide preview for planner reflection — same
            # privacy + vision gate as the page-mode screenshot (the preview
            # shows the data).
            _slides_preview_b64 = self._first_preview_base64(preview_images)
            if _slides_preview_b64 and allow_llm_see_data and _model and getattr(_model, "supports_vision", False):
                observation["summary"] += " First slide preview is attached — review it for visual correctness."
                observation["images"] = [{
                    "data": _slides_preview_b64,
                    "media_type": "image/png",
                    "source_type": "base64",
                }]

        if warnings:
            observation["warnings"] = warnings

        yield ToolEndEvent(
            type="tool.end",
            payload={
                "output": output,
                "observation": observation,
            }
        )

    def _trim_none(self, obj: Any) -> Any:
        """Remove None values and empty collections from nested structures."""
        try:
            if isinstance(obj, dict):
                out = {}
                for k, v in obj.items():
                    tv = self._trim_none(v)
                    if tv is None:
                        continue
                    if isinstance(tv, (dict, list)) and len(tv) == 0:
                        continue
                    out[k] = tv
                return out
            if isinstance(obj, list):
                items = [self._trim_none(v) for v in obj]
                return [v for v in items if not (v is None or (isinstance(v, (dict, list)) and len(v) == 0))]
            return obj
        except Exception:
            return obj

    def _build_slides_prompt(
        self,
        user_prompt: str,
        title: str | None,
        viz_profiles: List[Dict[str, Any]],
        instructions_context: str,
        report_title: str | None,
        allow_llm_see_data: bool,
        messages_context: str = "",
        image_count: int = 0,
        organization_settings: Any = None,
        files: Optional[List[Dict[str, Any]]] = None,
    ) -> str:
        """Build the prompt for generating slides using python-pptx code."""
        viz_json = json.dumps(viz_profiles, indent=2, default=str)

        # Embeddable art. Only images: python-pptx cannot place a PDF, and the
        # executor only loads image/* bytes, so advertising anything else would
        # promise an id that image() will reject.
        embeddable = [
            f for f in (files or [])
            if str(f.get("content_type") or "").startswith("image/")
        ]
        if embeddable:
            listed = "\n".join(
                f"  - {f['id']} — {f.get('filename') or 'image'}" for f in embeddable
            )
            embeddable_images = f"\n\n  Available image ids:\n{listed}"
        else:
            embeddable_images = (
                "\n\n  No images are attached to this deck — image_ids is empty. "
                "Build the design from shapes, color and type."
            )

        language_directive = build_language_directive(organization_settings)

        # Build attached images context
        images_context = ""
        if image_count > 0:
            images_context = f"\n**Attached Images:** {image_count} image(s) provided for visual reference. Use these to understand the design intent, branding, color schemes, or layout preferences the user wants to incorporate."

        return f"""Role: presentation author using python-pptx.{language_directive}
Generate python-pptx code to create a polished slide deck.

═══════════════════════════════════════════════════════════════════════════════
DECK CRAFT — decide this BEFORE writing any code
═══════════════════════════════════════════════════════════════════════════════

**1. Settle the storyline first.** Write the argument in one sentence: what
should the audience believe or do after seeing this? Every slide either
supports that sentence or gets cut. Then order slides so each earns the next.
The arc that works for analytical decks:

  1. Title — subject, audience, date.
  2. The headline — the single most important finding, stated outright. Do not
     save the conclusion for the end; executives read the first two slides.
  3. Evidence — one slide per supporting point, each with a chart.
  4. What changed / why — drivers, segments, root cause.
  5. So what — implications, risks, recommended actions.
  6. Appendix — detail, methodology, caveats.

  For a status or review deck, replace 2-5 with: where we are → what moved →
  what's blocked → what's next.

**2. Titles carry the message.** The title is the one line everyone reads.
Make it the finding, not the topic:
  - Weak:   "Revenue by Region"
  - Strong: "EMEA drove all of Q3 growth; every other region was flat"
A reader should page through titles alone and get the whole argument. If a
title could sit on any deck in any quarter, it is a label, not a takeaway.
Keep titles under ~12 words so they fit one line.

**3. One idea per slide.** If a slide needs "and" twice to explain, split it.
  - One chart per slide, unless two are being directly compared.
  - At most 5 bullets, at most 2 lines each, no sub-bullets.
  - No paragraphs. If prose is needed, the deliverable is a document, not a deck.
  - Numbers carry units and periods ("$4.2M, Q3" — not "4200000").
  - Put supporting detail in speaker notes via `slide.notes_slide`, not in
    shrunken body text.

**4. Never invent a number.** Every figure in a title or takeaway must match
what the chart shows. If the data does not support the claim, change the claim.

**5. Hold ONE visual system across the whole deck.** Pick a palette and stick
to it for every slide — same background family, same accent, same type scale,
same margins. A deck where slide 3 is light and slide 5 is dark, or where card
colors change without meaning, reads as broken no matter how good any single
slide is.

**6. A deck does NOT require data.** `visualizations` is often EMPTY — a topic,
narrative or announcement deck ("a deck about the 2026 World Cup") has no
charts at all, and that is a valid deck, not an error.

  - **Never index `visualizations[0]` without checking the list first.** On an
    empty list that raises IndexError and loses the whole deck. Guard every
    data-driven slide with `if visualizations:` and skip it otherwise.
  - With no data, carry the design with type, color, shapes and images: a
    full-bleed title, a section divider, a numbered-point layout, a quote, a
    stat stated as large type (only if the user supplied the number).
  - Do NOT invent charts, metrics or figures to fill the space. A confident
    typographic slide beats a fabricated bar chart.
  - Requested slide count is a hard constraint: "2 slides" means exactly 2.

═══════════════════════════════════════════════════════════════════════════════
AVAILABLE IN NAMESPACE (already provided — do not import)
═══════════════════════════════════════════════════════════════════════════════

Python-pptx classes and functions:
- Presentation, Inches, Pt, Emu, RGBColor
- PP_ALIGN, MSO_ANCHOR, MSO_SHAPE
- XL_CHART_TYPE, XL_LEGEND_POSITION
- CategoryChartData, ChartData

Note: Inches, Pt, Emu are functions, not methods.
   Use: Inches(1), Pt(24), Emu(914400)
   Not: 1.inches, 24.pt, value.inches

Data variables:
- visualizations: List[Dict] — each has 'title', 'columns', 'rows'
- report: Dict with 'id', 'title', 'theme'

Images (only present when the user attached or generated some):
- image_ids: List[str] — the embeddable image ids available to this deck
- image(file_id) -> stream — pass straight to add_picture. Returns a fresh
  stream per call, so the same image may be placed on several slides:
    pic = slide.shapes.add_picture(image(image_ids[0]), Inches(0), Inches(0),
                                   width=Inches(13.333))
  Cover the slide for a hero/background, or inset it in a content column.
  There is no filesystem access — `image()` is the only way to place art.
  When an image sits behind text, draw a translucent scrim rectangle between
  them or the text becomes unreadable; send the picture to the back by
  inserting it first.{embeddable_images}

Output:
- _pptx_output_path: str — path to save the presentation to

═══════════════════════════════════════════════════════════════════════════════
YOUR VISUALIZATIONS
═══════════════════════════════════════════════════════════════════════════════

{viz_json}

{"(Full sample data included above)" if allow_llm_see_data else "(Data samples hidden for privacy - use column names and row_count)"}

═══════════════════════════════════════════════════════════════════════════════
TASK
═══════════════════════════════════════════════════════════════════════════════

**Report Title:** {report_title or title or 'Presentation'}
**User Request:** {user_prompt}
{images_context}
{f"**Organization Instructions:** {instructions_context}" if instructions_context else ""}

═══════════════════════════════════════════════════════════════════════════════
PYTHON-PPTX QUICK REFERENCE
═══════════════════════════════════════════════════════════════════════════════

**Setup (16:9 widescreen):**
```python
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
```

**Add blank slide with dark background:**
```python
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(15, 23, 42)  # slate-900
```

**Add text box:**
```python
txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Title Text"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER
```

**Add bar chart (use this pattern for charts):**
```python
chart_data = CategoryChartData()
chart_data.categories = ['Q1', 'Q2', 'Q3', 'Q4']
chart_data.add_series('Revenue', (1.2, 1.5, 1.8, 2.1))

x, y, cx, cy = Inches(1), Inches(2), Inches(11), Inches(5)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data
).chart

# Style the chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False
plot = chart.plots[0]
plot.has_data_labels = True
```

**Other chart types:**
- XL_CHART_TYPE.COLUMN_CLUSTERED - vertical bars
- XL_CHART_TYPE.LINE - line chart
- XL_CHART_TYPE.PIE - pie chart
- XL_CHART_TYPE.AREA - area chart

**Dark background (slate-900 = RGB(15, 23, 42)):**
```python
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(15, 23, 42)
```

**Access visualization data** (only inside an `if visualizations:` guard — the
list is empty for a narrative deck):
```python
viz = visualizations[0]
# Each entry of viz['columns'] is a DICT, e.g. {{'field': 'Revenue', 'headerName': 'Revenue'}}.
# Extract the row keys from the 'field' values first:
fields = [c['field'] if isinstance(c, dict) else str(c) for c in viz['columns']]  # e.g. ['AlbumTitle', 'Revenue']
rows = viz['rows']        # list of dicts like {{'AlbumTitle': 'Greatest Hits', 'Revenue': 1500.0}}

# Get categories and values for a chart:
categories = [str(row.get(fields[0], '')) for row in rows]  # First column as labels
values = [float(row.get(fields[1]) or 0) for row in rows]   # Second column as values

# IMPORTANT: never pass a column dict itself to a text property or a row index —
# use fields[i] (a string like 'Revenue'): row[fields[1]] is row['Revenue']
```

═══════════════════════════════════════════════════════════════════════════════
DESIGN PHILOSOPHY - CREATE BEAUTIFUL, PROFESSIONAL SLIDES
═══════════════════════════════════════════════════════════════════════════════

**COLOR STRATEGY - Be Topic-Specific:**
Choose colors that feel designed for THIS topic. If your colors would work for any presentation, you haven't made specific enough choices.

Structure: One DOMINANT color (60-70% visual weight), 1-2 supporting tones, one accent.

Example palettes (pick one that fits the topic):
- **Midnight Executive**: Navy (0,31,63), Steel (119,136,153), Gold accent (212,175,55)
- **Forest & Moss**: Deep green (34,87,76), Sage (138,154,91), Cream (245,245,220)
- **Coral Energy**: Coral (255,127,80), Teal (0,128,128), Sand (244,232,214)
- **Ocean Depths**: Deep blue (0,51,102), Aqua (0,180,180), Pearl (240,248,255)
- **Sunset Warm**: Burgundy (128,0,32), Orange (255,140,0), Cream (255,253,240)
- **Modern Minimal**: Charcoal (54,69,79), Light gray (220,220,220), Teal accent (0,150,136)

**Layout variety — vary between slides:**
Every slide should have visual elements — charts, shapes, or decorative elements. Avoid text-only slides.

Vary layouts between:
- Two-column (text left, chart right or vice versa)
- Full-width chart with title above
- KPI cards in a row (3-4 metric boxes)
- Chart with callout boxes for key insights
- Split layout with accent shape dividers

**Typography:**
- Titles: 36-44pt bold, interesting positioning (not always centered)
- Body text: 18-24pt, left-aligned (avoid center-aligning body text)
- KPI numbers: 48-72pt bold for impact
- Use font color contrast: white on dark, dark on light accents

**VISUAL ELEMENTS TO ADD:**
- Accent shapes: rectangles, rounded rectangles for backgrounds
- Divider lines or shapes between sections
- Colored boxes behind KPI numbers
- Subtle shape overlays for visual interest

**Common mistakes to avoid:**
- Using `value.inches` instead of `Inches(value)` — Inches/Pt/Emu are functions.
- Repeating the same layout across slides — vary it.
- Center-aligning body text — use left alignment.
- Using only blue without topic-specific reasoning.
- Text-only slides without visual elements.
- Accent lines directly under titles (hallmark of generic slides).
- Cramming too much data — limit charts to top 8-10 items.
- Adding a chart without checking its rows first — `CategoryChartData` raises
  "chart data contains no categories" on an empty list and that failure loses
  the whole deck, not just the slide.

**Technical requirements:**
1. Define `generate_slides(visualizations, report)` returning a Presentation.
2. Use 16:9 widescreen: Inches(13.333) x Inches(7.5).
3. Create real charts with slide.shapes.add_chart() + CategoryChartData.
4. Use visualization data from the visualizations list. Read rows with
   `viz.get('rows', [])` and ALWAYS guard before charting:
   `if rows:` — build the chart; otherwise render the slide without it.
5. Margins: start shapes at Inches(0.75) to Inches(1) from edges.

**Rendering defects to prevent (these are what make a deck look broken):**
- **Charts on dark backgrounds render unreadable.** python-pptx defaults every
  chart label to near-black, which disappears on a dark slide. On a dark
  background set them explicitly — category and value tick labels, data
  labels, and the chart title:
  `chart.font.color.rgb = LIGHT` plus
  `chart.category_axis.tick_labels.font.color.rgb = LIGHT` and
  `chart.value_axis.tick_labels.font.color.rgb = LIGHT`.
  Pick series colors that contrast with the background too.
- **Content running off the slide.** The canvas is 7.5in tall. Keep every shape
  between Inches(0.4) and Inches(7.1) vertically, and inside Inches(0.75) from
  the left and right edges — a title at left=0 reads as a layout bug. Before
  emitting a card grid, check `top + height` for the LAST row fits.
- **Titles colliding with what follows.** A long title wraps to 2-3 lines. Give
  the title box enough height for the wrap and start the next element BELOW it;
  never overlap a subtitle, accent line or chart with the title block.

═══════════════════════════════════════════════════════════════════════════════
OUTPUT FORMAT - Example with Design Principles Applied
═══════════════════════════════════════════════════════════════════════════════

```python
def generate_slides(visualizations, report):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color palette - choose colors that fit the topic
    PRIMARY = RGBColor(0, 51, 102)      # Deep blue
    SECONDARY = RGBColor(0, 128, 128)   # Teal
    ACCENT = RGBColor(255, 140, 0)      # Orange accent
    BG_DARK = RGBColor(15, 23, 42)      # Dark background
    TEXT_LIGHT = RGBColor(255, 255, 255)
    TEXT_MUTED = RGBColor(148, 163, 184)

    def set_background(slide, color=BG_DARK):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_accent_shape(slide, left, top, width, height, color):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 1: Title with accent shape
    # ═══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    # Accent shape behind title
    add_accent_shape(slide, Inches(0), Inches(2.5), Inches(5), Inches(2.5), PRIMARY)

    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(3), Inches(12), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = report.get('title', 'Presentation')
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = TEXT_LIGHT

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 2: KPI Cards Row (if we have numeric data)
    # ═══════════════════════════════════════════════════════════════
    if visualizations and visualizations[0].get('rows'):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_background(slide)

        viz = visualizations[0]
        rows = viz.get('rows', [])
        columns = viz.get('columns', [])

        # Create 3 KPI cards across the slide
        card_width = Inches(3.5)
        card_height = Inches(2.5)
        start_x = Inches(1)
        card_y = Inches(2.5)
        gap = Inches(0.5)

        for i, col in enumerate(columns[:3]):
            if i >= 3:
                break
            x = start_x + i * (card_width + gap)

            # Card background
            card = add_accent_shape(slide, x, card_y, card_width, card_height, PRIMARY)

            # Value (large number)
            val = rows[0].get(col, 0) if rows else 0
            val_box = slide.shapes.add_textbox(x + Inches(0.3), card_y + Inches(0.5), card_width - Inches(0.6), Inches(1.2))
            tf = val_box.text_frame
            p = tf.paragraphs[0]
            p.text = "{{:,.0f}}".format(float(val)) if isinstance(val, (int, float)) else str(val)
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = TEXT_LIGHT

            # Label
            label_box = slide.shapes.add_textbox(x + Inches(0.3), card_y + Inches(1.7), card_width - Inches(0.6), Inches(0.6))
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = col
            p.font.size = Pt(14)
            p.font.color.rgb = TEXT_MUTED

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 3: Chart with title (different layout)
    # ═══════════════════════════════════════════════════════════════
    if visualizations:
        viz = visualizations[0]
        columns = viz.get('columns', [])
        rows = viz.get('rows', [])

        if len(columns) >= 2 and rows:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            set_background(slide)

            # Title on left side
            title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(5), Inches(1))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = viz.get('title', 'Data Analysis')
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = TEXT_LIGHT

            # Extract data
            col_label = columns[0]
            col_value = columns[1]
            categories = [str(row.get(col_label, ''))[:20] for row in rows[:8]]
            values = [float(row.get(col_value, 0) or 0) for row in rows[:8]]

            # Chart (full width below title)
            chart_data = CategoryChartData()
            chart_data.categories = categories
            chart_data.add_series(col_value, tuple(values))

            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.BAR_CLUSTERED,
                Inches(0.75), Inches(1.5), Inches(11.833), Inches(5.5),
                chart_data
            ).chart
            chart.has_legend = False

    return prs

# Execute and save
prs = generate_slides(visualizations, report)
prs.save(_pptx_output_path)
```

Create a beautiful, varied presentation following these design principles. Each slide should look DIFFERENT from the others. Use visual elements, accent shapes, and thoughtful color choices:"""

    def _build_page_system_prompt(self) -> str:
        """Static system prompt for page/dashboard generation.

        Contains only stable reference material (sandbox runtime docs,
        component contract, filtering rules, design guidance, output format).
        Kept free of per-call state so provider prompt caching can reuse it
        across every create/edit call.
        """
        return f"""Role: frontend developer and data visualization engineer. You build React dashboard artifacts from the user's design request and visualization data (both provided in the user message), following this reference.

═══════════════════════════════════════════════════════════════════════════════
REFERENCE — TOOLS, COMPONENTS & DATA
═══════════════════════════════════════════════════════════════════════════════

{SANDBOX_RUNTIME_PROMPT}

CHARTING:

**`<EChart height={{N}} option={{{{...}}}} />`** — chart wrapper. Supports ALL ECharts chart types. 'bow' theme pre-configures colors, tooltip, grid, axes. For standard charts, only write data mapping:
```jsx
<EChart height={{300}} option={{{{ xAxis: {{ type: 'category', data: rows.map(r => r.name) }}, yAxis: {{ type: 'value' }}, series: [{{ type: 'bar', data: rows.map(r => r.val) }}] }}}} />
<EChart height={{300}} option={{{{ tooltip: {{ trigger: 'item' }}, series: [{{ type: 'pie', radius: ['45%','75%'], data: rows.map(r => ({{ value: r.amt, name: r.lbl }})) }}] }}}} />
<EChart height={{300}} option={{{{ xAxis: {{ type: 'category', data: rows.map(r => r.date) }}, yAxis: {{ type: 'value' }}, series: [{{ type: 'line', data: rows.map(r => r.val), areaStyle: {{ opacity: 0.15 }} }}] }}}} />
```
For advanced charts (radar, gauge, treemap, sunburst, funnel, sankey, calendar heatmap, parallel coordinates, graph), pass the full ECharts option — the theme still provides colors and tooltip:
```jsx
<EChart height={{300}} option={{{{ radar: {{ indicator: indicators }}, series: [{{ type: 'radar', data: radarData }}] }}}} />
<EChart height={{250}} option={{{{ series: [{{ type: 'gauge', data: [{{ value: 72 }}], detail: {{ formatter: '{{value}}%' }} }}] }}}} />
<EChart height={{400}} option={{{{ series: [{{ type: 'treemap', data: treeData }}] }}}} />
```

AVAILABLE COMPONENTS (convenience shortcuts — not requirements):
- `<KPICard title="" value={{fmt(n, {{currency:true}})}} subtitle="" color="#3B82F6" className="" titleClassName="" subtitleClassName="" style={{{{}}}} />` — `className` replaces default theme (bg-white, border, text-slate-900). `titleClassName`/`subtitleClassName` replace title/subtitle defaults. `style` for inline overrides. Theme these to match your color story:
  - Dark: `className="bg-slate-900 border-slate-700 text-white" titleClassName="text-slate-400"`
  - Colored: `className="bg-indigo-50 border-indigo-200 text-indigo-900" titleClassName="text-indigo-600"`
- `<SectionCard title="" subtitle="" className="" titleClassName="" subtitleClassName="" style={{{{}}}}>...children...</SectionCard>` — same theming: `className` replaces defaults, `titleClassName`/`subtitleClassName` for text. Theme to match.
- `<FilterSelect label="" options={{arr}} selected={{arr}} onChange={{fn}} searchable={{bool}} className="" style={{{{}}}} />` — multi-select dropdown (portaled). Built-in search at 8+ options. `className` replaces default theme (bg-white border-slate-200 text-slate-900) — pass e.g. `className="bg-slate-900 border-slate-700 text-slate-100"` for dark.
- `<FilterSearch label="" value={{str}} onChange={{e => setFilter(field, e.target.value)}} placeholder="Search..." className="" style={{{{}}}} />` — text search. `className` replaces default theme.
- `<FilterDateRange label="" value={{filters[field] || {{}}}} onChange={{val => setFilter(field, val)}} type="date" className="" style={{{{}}}} />` — date range picker. `className` replaces default theme.
- `fmt(n, opts)` — `{{currency:true}}`, `{{pct:true}}`, auto K/M/B
- `<LoadingSpinner size={{32}} />`

All components are fully themeable via `className`/`titleClassName`/`subtitleClassName`/`style`. Don't leave default white/slate styling when your design calls for something different. If the design needs something these can't express — build custom React + Tailwind.

**HOST DARK MODE:** The sandbox runs Tailwind with `darkMode: 'class'`; the host toggles a `dark` class on `<html>` to match the viewer's theme (live, no reload). Component DEFAULTS already adapt (they carry `dark:` variants), as do the iframe body and charts (bow/bow-dark ECharts themes). So: prefer the defaults when no specific color story is requested — they look right in both modes. When you hardcode light colors on custom markup or via `className` overrides, pair them with `dark:` variants (`bg-white dark:bg-slate-900 text-slate-900 dark:text-slate-100`) unless the design is intentionally single-look (then set explicit backgrounds everywhere).

**INFO POPOVER (required):** Pass `viz={{viz[N]}}` to every `<KPICard>` and `<SectionCard>` you build from a visualization. This renders a small built-in "ⓘ" button that lets users inspect the data behind each component (Data tab with rows, Code tab with the query). Use the index of the visualization the card is derived from (the primary one if it combines several). When a card renders FILTERED rows (you called `filterRows(viz[N].rows)`), ALSO pass `rows={{<those filtered rows>}}` so the popover shows the filtered view that matches the component, not the full dataset. When a card AGGREGATES or derives its value client-side, ALSO pass `calc="<formula>"` describing the math with real column names, e.g. `calc="SUM(UnitPrice × Quantity) grouped by GenreName"` or `calc="COUNT(DISTINCT CustomerId)"` — the popover shows it as a "Calculation" line. If you render a chart with a bare `<EChart>` that is NOT inside a `<SectionCard>`, pass `viz={{viz[N]}}` (and `rows`/`calc` if relevant) to the `<EChart>` itself so it still gets the popover.

**CUSTOM MARKUP — add `data-bow-*` attributes (required):** Whenever you build your OWN containers instead of `<KPICard>`/`<SectionCard>`/`<EChart>` (custom `<div>` KPI tiles, chart wrappers, tables), annotate each item's outer element with `data-bow-viz="N"` (source visualization index) and `data-bow-calc="<formula>"` when the value is derived. A global overlay then renders the same Data/Code/Calc popover on each item. Example: `<div data-bow-viz={{0}} data-bow-calc="SUM(UnitPrice × Quantity)">...custom tile...</div>`. EVERY metric, chart, and table must be reachable via either a prebuilt component's `viz` prop OR a `data-bow-viz` attribute — never leave an item with no way to inspect its data.

DATA ACCESS:

```javascript
const data = useArtifactData(); // Returns null while loading
// data = {{ report: {{id, title}}, visualizations: [...] }}
```

Each visualization:
```js
{{
  id: "uuid",
  title: "Visualization Title",
  columns: [{{ "headerName": "Album Title", "field": "AlbumTitle", "dtype": "object", "unique_count": 150 }}, ...],
  rows: [{{ "AlbumTitle": "Battlestar Galactica", "total_revenue": 35.82 }}, ...],
  row_count: 4321,   // TRUE dataset size
  view: {{ /* chart config hints */ }},
  dataModel: {{ /* series/axis config */ }}
}}
```

- Use `column.field` to access row values: `row[column.field]`
- Use `column.headerName` for display labels
- Column metadata includes `dtype` (pandas type) and `unique_count` — use these for filter/format decisions
- **Do not hardcode data** — all values should come from `data.visualizations[N].rows`
- **Sample vs full data:** the `rows`/`sample_rows` shown in the user message are a SAMPLE (capped at 100 rows per visualization) for generation and preview. At runtime the dashboard receives the FULL dataset — `row_count` rows. `row_count` is the true dataset size; `sample_row_count` is the sample size. Write code that works on the full dataset (aggregate with reduce/Map, paginate long tables) and NEVER hardcode workarounds for the sample size.
- **Defensive coding**: Row values and properties can be `null`/`undefined`. Use optional chaining or fallbacks before calling `.includes()`, `.toLowerCase()`, `.startsWith()`, `.split()`, etc. Example: `(row.name || '').includes('x')` or `String(val ?? '').toLowerCase()`. Do not call string methods on a value that could be nullish.

View hints — honor the viz config:
The `view_config` on each visualization describes how the author wants the data rendered. Follow it when generating code.

- `view_config.aggregation` (`"sum" | "avg" | "count" | "min" | "max"`): the raw rows are granular, so aggregate the relevant value column before rendering (especially for `count`, `metric_card`, `pie_chart`, `heatmap`). Use `rows.reduce(...)`. Example for a metric card with aggregation=sum:
  ```js
  const total = useMemo(
    () => viz[0].rows.reduce((s, r) => s + (Number(r.revenue) || 0), 0),
    [viz]
  );
  ```
  For pie/heatmap/bar charts that group by a category, group first and aggregate the value per group rather than using the first matching row.

- `view_config.series_aggregations` (array of `{{key, aggregation}}`): apply the given aggregation per series when building multi-series bar/line/area charts.

- `view_config.default_filters` (array of `{{column, operator, value}}`): the author wants the dashboard to open with these filters already applied. Seed them on first mount so the initial view matches the intent, for example:
  ```js
  const {{ filters, setFilter, filterRows }} = useFilters();
  useEffect(() => {{
    // Seed defaults once — operators follow the useFilters contract.
    {{/* for each entry in view_config.default_filters */}}
    setFilter('column_name', value);
  }}, []);
  ```
  If the underlying runtime uses richer operators (`equals`, `greater_than`, etc.), either call `setFilter` with the operator-aware object it expects, or compute the filtered rows directly via `filterRows(viz[N].rows)` once the filter is seeded. Render the filtered view when defaults are present so the initial numbers match the author's intent.

FILTERING:
- Use `useFilters()` hook for cross-visualization filtering — returns `{{ filters, setFilter, resetFilters, filterRows }}`
- YOU choose which columns to filter — use `dtype` and `unique_count` from the column metadata:
  - `<FilterSelect>` for low-cardinality columns (`unique_count` < ~50, dtype "object"/"int64" with few values)
  - `<FilterSearch>` for high-cardinality text columns (`unique_count` > 50, dtype "object")
  - `<FilterDateRange>` for date/time columns (dtype contains "datetime" or values are date strings)
- Get unique values directly: `[...new Set(viz[N].rows.map(r => r[field]))]`

FILTER FEASIBILITY AUDIT — DO THIS FIRST, BEFORE WRITING CODE:
Before wiring any cross-viz filter, verify it will actually work. A filter that looks wired but silently leaves some vizs untouched is a broken dashboard, not a partial one.

For each dimension you intend to filter by:
1. **Enumerate participating vizs** — which vizs should this filter affect? (Usually: any viz whose topic logically shares the dimension, e.g. a "customer" filter should affect every viz about customers, payments, orders, etc.)
2. **Check column presence** — does each participating viz have the filter column (directly, or via a rename you can handle with `fieldMap`)? Check the `columns` array in YOUR VISUALIZATIONS below.
3. **Decide per dimension**:
   - ALL participants have the column → wire the global filter, use `fieldMap` for renames.
   - SOME participants lack the column but the gap is genuine (no join key in the source data) → make the filter LOCAL to the vizs that support it; do not pretend it affects others.
   - SOME participants lack the column but they should have it (the underlying data supports it, the query just didn't project the column) → **do not wire the filter; do not build the dashboard with a dead filter.** End your response by reporting the gap so the planner can recreate the offending queries before you try again. Example: "Cannot wire `customer_id` filter — `payments` viz lacks `customer_id` but `payments.customer_id` exists in schema. Recreate the payments query with `customer_id` projected, then retry create_artifact."

FILTER PLACEMENT — global vs local:
- **Global filter** (column present in 2+ vizs AFTER the audit above): place in a top-level filter bar above all content. Use one shared filter + `fieldMap` for renames, not duplicates.
- **Local filter** (column present in only 1 viz): place INSIDE that viz's `<SectionCard>`, visually next to the chart/table it affects.
- When a filter affects multiple vizs, add visible UI indication that they're linked.

FILTER DATA FLOW:
- Every viz that passes the feasibility audit for a filter should use `filterRows()` as its data source — for charts, tables, and any KPI/summary derived from that viz.
- KPI cards that summarize filtered data (sum, count, avg) should be computed from filtered rows, not from raw `viz[N].rows`.
- Do not call `filterRows` on a viz that doesn't have the filter column just to "be safe" — silently passing rows through makes the filter look active when it isn't. Audit first, wire second.

EXAMPLE 1 — Global "region" filter affecting KPIs + bar chart + table:
  const {{ filters, setFilter, resetFilters, filterRows }} = useFilters();
  const regions = useMemo(() => [...new Set(vizSales.rows.map(r => r.region))], [vizSales]);
  // ALL downstream from vizSales uses filtered:
  const filteredSales = filterRows(vizSales.rows);
  const totalRevenue = useMemo(() => filteredSales.reduce((s, r) => s + r.revenue, 0), [filteredSales]);
  const chartData = useMemo(() => ({{ labels: filteredSales.map(r => r.month), values: filteredSales.map(r => r.revenue) }}), [filteredSales]);
  // Cross-viz filtering with field mapping:
  const filteredDetails = filterRows(vizDetails.rows, {{ region: 'RegionName' }});
  // Layout: <FilterSelect> in top bar, KPIs below, charts below that

EXAMPLE 2 — Local filter inside a SectionCard:
  const {{ filters, setFilter, filterRows }} = useFilters();
  const filtered = filterRows(vizProducts.rows);
  // Layout: <SectionCard title="Products"><FilterSelect .../><EChart ... /></SectionCard>

- Include a Reset button when any filters are active (`Object.keys(filters).length > 0`)
- After filtering, if a visualization has zero matching rows, display "No data matches current filters"

═══════════════════════════════════════════════════════════════════════════════
DESIGN GUIDANCE (use when the user hasn't specified styling)
═══════════════════════════════════════════════════════════════════════════════

If the user specified a theme/style/colors above, follow that — skip this section.
Otherwise, design a visually striking, publication-quality dashboard — not a generic template.

COLOR & IDENTITY:
- Pick a cohesive color story that fits the data topic. A finance dashboard should feel different from a music dashboard, which should feel different from a healthcare dashboard.
- Choose one dominant color (60-70%), 1-2 supporting tones, and one accent for highlights/CTAs.
- Do NOT default to generic blue. Blue is fine if it fits the topic — but earn it, don't default to it.
- Theme ALL components (KPICard, SectionCard, filters) to match — use `className`, `titleClassName`, `subtitleClassName` props. Default white/slate is only appropriate for a clean/minimal design intent.

LAYOUT & HIERARCHY:
- Lead with the most important insight — KPIs or headline metric at the top.
- Create clear visual hierarchy: primary chart large, secondary charts smaller, supporting data compact.
- Use intentional whitespace — not "fill every pixel" but not "float in empty space" either.
- Vary card sizes and chart heights to create rhythm. A grid of same-sized boxes is boring.

TYPOGRAPHY & POLISH:
- Clean, modern typography. Titles concise and descriptive, not generic ("Revenue by Region" not "Chart 1").
- Subtle shadows, rounded corners, light borders — enough depth to feel crafted, not flat.
- Light mode default. Dark mode only if the topic or user suggests it.

CHART SELECTION:
- Choose the best visualization for the data shape — don't default to bar charts for everything.
- Standard charts (bar, line, pie, area) for simple relationships. Advanced charts (radar, gauge, treemap, funnel, sankey, heatmap) when the data structure rewards it.
- Show data from different angles without redundancy. Each chart should reveal something the others don't.

The goal: it should look like a designer built it for this specific dataset, not like a template was filled in.

═══════════════════════════════════════════════════════════════════════════════
RESPONSIVE LAYOUT (REQUIRED — always applies, even when the user specified a theme/style)
═══════════════════════════════════════════════════════════════════════════════

The dashboard is embedded in an iframe whose width is NOT fixed — the SAME code renders in a narrow chat side-panel (~360–480px), a normal report view (~900px), and a full-screen / published view (up to ~1920px). It MUST reflow gracefully at every width with NO horizontal page scroll and NO clipped or squished content. Build it fluid and mobile-first; only deviate if the user EXPLICITLY asked for a fixed width.

Concrete rules — follow all of them:
- **Outer container:** fluid width, never a fixed pixel width. Use `w-full min-h-full` (add `max-w-screen-2xl mx-auto` only if you want to cap width on huge screens). Responsive padding: `p-4 md:p-6 lg:p-8`. NEVER `w-[1200px]`, `min-w-[...]`, or any fixed-pixel width on layout containers.
- **KPI / stat rows:** use a responsive grid that collapses on narrow screens, e.g. `grid grid-cols-2 md:grid-cols-4 gap-4` (2-up on mobile → 4-up on desktop). Do NOT use a flex row of fixed-width cards that overflows.
- **Chart grids:** start single-column and add columns at breakpoints, e.g. `grid grid-cols-1 lg:grid-cols-2 gap-6`. A primary/feature chart can stay full-width (`col-span-full` or its own row). Never lock a multi-column grid with no single-column fallback.
- **Charts:** give each `<EChart>` a `w-full` container and a fixed `height` (px) — it auto-resizes to its container via ResizeObserver, so width takes care of itself. Do not set a pixel width on charts.
- **Tables & wide content:** wrap in `<div className="overflow-x-auto">` so a wide table scrolls inside its card instead of blowing out the page width. Use `min-w-full` on the `<table>`, not a fixed width.
- **Filter bars:** `flex flex-wrap gap-3` so filters wrap to the next line on narrow widths instead of overflowing.
- **Text & numbers:** allow large KPI numbers to scale (e.g. `text-2xl md:text-3xl`) and use `truncate`/`break-words` where labels can be long, so nothing overflows its card.
- **Sanity check:** before finishing, mentally render at ~380px wide — every row must wrap to 1–2 columns, no element wider than the viewport, no horizontal scrollbar on the body.

═══════════════════════════════════════════════════════════════════════════════
OUTPUT FORMAT
═══════════════════════════════════════════════════════════════════════════════

```
<script type="text/babel">
function App() {{
  const data = useArtifactData();
  if (!data) return <div className="flex items-center justify-center h-screen text-gray-400"><LoadingSpinner size={{32}} /></div>;
  const viz = data.visualizations;
  // ... concise dashboard code
}}
ReactDOM.createRoot(document.getElementById('root')).render(<App />);
</script>
```

Structure: all code should be inside `function App() {{ ... }}` with `ReactDOM.createRoot(document.getElementById('root')).render(<App />);` at the end. Do not put return statements outside a function.

Rules: `<script type="text/babel">` wrapper. `useArtifactData()` for data. `<EChart option={{...}} />` for charts. Pass `viz={{viz[N]}}` to every KPICard/SectionCard so the built-in info popover shows the data behind it. RESPONSIVE — fluid width, responsive grids (`grid-cols-1 md:grid-cols-2 lg:grid-cols-N`), no fixed-pixel widths, no horizontal page scroll at any width (see RESPONSIVE LAYOUT section above); required unless the user asked for a fixed width. Handle zero rows. No hardcoded data. No UUIDs/branding/emoji. Guard nullish values before string methods (use `(val || '')` or `String(val ?? '')`).

**Code size:** Write compact code — no unnecessary variables, comments, or verbose JSX. Omit default props. Don't repeat theme styling the 'bow' theme already provides. Prefer inline expressions over separate variables when used once. For simple dashboards target under 8K characters. For detailed/specific user requests, use as much space as needed to faithfully implement their design — fidelity to the user's request is more important than brevity."""

    def _build_page_prompt(
        self,
        user_prompt: str,
        title: str | None,
        viz_profiles: List[Dict[str, Any]],
        instructions_context: str,
        report_title: str | None,
        allow_llm_see_data: bool,
        messages_context: str = "",
        image_count: int = 0,
        organization_settings: Any = None,
        files: Optional[List[Dict[str, Any]]] = None,
        identity_context: str = "",
    ) -> str:
        """Build the dynamic user prompt for page/dashboard generation.

        Static reference material lives in _build_page_system_prompt (cached
        as the system prompt); this carries only per-call state.
        """
        viz_json = json.dumps(viz_profiles, indent=2, default=str)

        # Server-side query parameters: when any viz declares them, the
        # dashboard must wire controls through useParams() (not useFilters).
        params_directive = ""
        _param_profiles = [
            (p.get("title") or p.get("id"), p.get("parameters"))
            for p in viz_profiles if p.get("parameters")
        ]
        if _param_profiles:
            _lines = "\n".join(
                f"  - {title}: " + ", ".join(
                    f"{d.get('name')} ({d.get('type', 'string')}, {d.get('source', 'input')})"
                    for d in decls
                )
                for title, decls in _param_profiles
            )
            params_directive = (
                "\n**Server-side query parameters (MANDATORY UI):** these visualizations' queries "
                "declare parameters — the platform re-runs them at the data source when a value "
                "changes (see useParams() in the runtime reference):\n" + _lines + "\n"
                "Rules: render ONE control per unique param name (same name across queries = one "
                "control driving all of them) wired to useParams().setParam(name, value); an "
                "optional param gets an 'All' choice that sets null; identity-source params get NO "
                "input — show a small 'scoped to you' badge; show a subtle loading state while "
                "useParams().loading is true, and ALWAYS render useParams().error when set (a "
                "small banner/toast) — a failed re-run must never look like a stale table. "
                "Do NOT emulate these with useFilters/client-side "
                "filtering — the fresh rows arrive through useArtifactData() automatically. "
                "A control's choice list must be STABLE: populate it from useParamOptions(name) "
                "(declared options / options-source query, host-resolved) — NEVER derive choices "
                "from the rows that control filters, or selecting a value collapses the list to "
                "the current selection. Bind option.value into setParam — never the label. "
                "A list-typed param renders as a MULTI-select (or checkable list) that submits "
                "an ARRAY of option.value entries; empty selection = null (All).\n"
            )

        language_directive = build_language_directive(organization_settings)

        # Build attached images context
        images_context = ""
        if image_count > 0:
            images_context = f"\n**Attached Images:** {image_count} image(s) provided for visual reference. Use these to understand the design intent, branding, color schemes, or layout preferences the user wants to incorporate."

        # Build embedded-files context (generated images / uploaded images or PDFs).
        # These are rendered via the <BowFile> sandbox component by file id.
        files_context = ""
        if files:
            lines = [
                f'- id="{f["id"]}"  type={f.get("content_type", "")}  name={f.get("filename", "")}'
                for f in files
            ]
            files_context = (
                "\n**Embedded Files (render with `<BowFile>`):** You MUST place each of these "
                "files in the layout using `<BowFile id=\"<id>\" />` (see the BowFile entry in the "
                "sandbox runtime docs). Images render inline; PDFs render in a viewer. Do NOT inline "
                "base64 and do NOT use a raw <img src>. Available files:\n" + "\n".join(lines)
            )

        # Note: Previous artifact code is now available via observation context (from create_artifact/read_artifact)
        # The planner can call read_artifact if needed to load previous code into context

        return f"""═══════════════════════════════════════════════════════════════════════════════
Design request (primary specification — takes precedence when it conflicts with reference defaults)
═══════════════════════════════════════════════════════════════════════════════

**Report Title:** {report_title or title or 'Dashboard'}
**User Request:** {user_prompt}
{images_context}
{files_context}
{params_directive}
{f"**Organization Instructions:**{chr(10)}{instructions_context}" if instructions_context else ""}
{identity_context}
{f"**Conversation History:**{chr(10)}{messages_context}" if messages_context else ""}
{language_directive}

If the user specified a theme, layout, colors, or style above — follow that exactly.
If the user did not specify styling, use the design guidance from the reference instructions.

YOUR VISUALIZATIONS:

{viz_json}

{"(Sample data included above — rows are a sample capped at 100; row_count is the true dataset size)" if allow_llm_see_data else "(Data samples hidden for privacy - use column names and row_count to understand the data structure)"}

Now create the dashboard:"""

    async def _build_file_datauris(self, db, included_files: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        """Read embedded files from disk and return them as data: URIs.

        Shared with the PDF export path — see _artifact_images.build_file_datauris.
        """
        from ._artifact_images import build_file_datauris

        return await build_file_datauris(db, included_files)

    def _build_prompt(
        self,
        user_prompt: str,
        title: str | None,
        mode: str,
        viz_profiles: List[Dict[str, Any]],
        instructions_context: str,
        report_title: str | None,
        allow_llm_see_data: bool,
        messages_context: str = "",
        image_count: int = 0,
        organization_settings: Any = None,
        files: Optional[List[Dict[str, Any]]] = None,
        identity_context: str = "",
    ) -> str:
        """Build the prompt for generating artifact code. Dispatches to mode-specific builders."""
        if mode == "slides":
            return self._build_slides_prompt(
                user_prompt=user_prompt,
                title=title,
                viz_profiles=viz_profiles,
                instructions_context=instructions_context,
                report_title=report_title,
                allow_llm_see_data=allow_llm_see_data,
                messages_context=messages_context,
                image_count=image_count,
                organization_settings=organization_settings,
                files=files,
            )
        return self._build_page_prompt(
            user_prompt=user_prompt,
            title=title,
            viz_profiles=viz_profiles,
            instructions_context=instructions_context,
            report_title=report_title,
            allow_llm_see_data=allow_llm_see_data,
            messages_context=messages_context,
            image_count=image_count,
            organization_settings=organization_settings,
            files=files,
            identity_context=identity_context,
        )

    def _extract_code(self, response: str, mode: str = "page") -> str:
        """Extract the code from the LLM response.

        For 'page' mode: Extract React code from <script type="text/babel"> tags
        For 'slides' mode: Extract python-pptx code from python code blocks
        """
        if mode == "slides":
            return self._extract_slides_python(response)

        # Dashboard mode - extract React code from script tags
        start_marker = "<script type=\"text/babel\">"
        end_marker = "</script>"

        start_idx = response.find(start_marker)
        if start_idx == -1:
            # Try alternative markers
            start_marker = "<script type='text/babel'>"
            start_idx = response.find(start_marker)

        if start_idx != -1:
            end_idx = response.find(end_marker, start_idx)
            if end_idx != -1:
                code = response[start_idx:end_idx + len(end_marker)]
                return self._sanitize_code(self._ensure_app_wrapper(code))

        # If no script tags found, wrap the response
        code = response.strip()
        if not code.startswith("<script"):
            code = f'<script type="text/babel">\n{code}\n</script>'

        return self._sanitize_code(self._ensure_app_wrapper(code))

    @staticmethod
    def _sanitize_code(code: str) -> str:
        """Fix common LLM code generation artifacts deterministically."""
        import re

        # Fix double-brace pattern: function App() {\n{ ... }\n}
        # The LLM sometimes wraps the function body in an extra block scope.
        # Match: function App() {\n{ at the start, and }\n} at the end (before render call)
        code = re.sub(
            r'(function\s+\w+\s*\([^)]*\)\s*\{)\s*\n\s*\{',
            r'\1',
            code,
        )
        # Remove the matching trailing extra }
        # Look for }\n}\n before ReactDOM.createRoot
        code = re.sub(
            r'\}\s*\n\s*\}\s*\n(\s*ReactDOM\.createRoot)',
            r'}\n\1',
            code,
        )

        return code

    @staticmethod
    def _ensure_app_wrapper(code: str) -> str:
        """Ensure code has a proper App component wrapper.

        LLM sometimes outputs bare return statements outside a function.
        Detect and fix by wrapping the inner code in function App() + ReactDOM.createRoot.
        """
        import re

        # Check if code already has an App function/component
        if re.search(r'function\s+App\s*\(', code) or re.search(r'(?:const|let|var)\s+App\s*=', code):
            return code

        # Extract inner code between script tags
        inner_match = re.search(
            r'<script\s+type=["\']text/babel["\']>\s*([\s\S]*?)\s*</script>',
            code
        )
        if not inner_match:
            return code

        inner = inner_match.group(1).strip()

        # Strip any existing broken ReactDOM.createRoot/render calls
        inner = re.sub(r'ReactDOM\.createRoot\(.*?\)\.render\(.*?\);?\s*$', '', inner, flags=re.DOTALL).strip()

        logger.warning("_ensure_app_wrapper: LLM output missing function App() wrapper — auto-wrapping")

        wrapped = (
            '<script type="text/babel">\n'
            'function App() {\n'
            f'{inner}\n'
            '}\n'
            "ReactDOM.createRoot(document.getElementById('root')).render(<App />);\n"
            '</script>'
        )
        return wrapped

    def _extract_slides_python(self, response: str) -> str:
        """Extract python-pptx code for slides mode."""
        import re

        # Try to find Python code block
        python_match = re.search(r'```python\s*([\s\S]*?)```', response)
        if python_match:
            return python_match.group(1).strip()

        # Try generic code block
        code_match = re.search(r'```\s*([\s\S]*?)```', response)
        if code_match:
            return code_match.group(1).strip()

        # Look for function definition as start marker
        func_start = response.find('def generate_slides')
        if func_start != -1:
            # Find the prs.save() call at the end
            save_end = response.rfind('prs.save(')
            if save_end != -1:
                # Include the full save line
                end_idx = response.find(')', save_end)
                if end_idx != -1:
                    return response[func_start:end_idx + 1].strip()
            return response[func_start:].strip()

        # Fallback: return the response as-is
        return response.strip()

