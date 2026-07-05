#!/usr/bin/env python
"""Generate a realistic 'network directory' of files to exercise the
`network_dir` data source connector (list / search / read / write).

Produces a nested tree under a target directory:

    <root>/
      contracts/            many CSVs + a few PDFs/markdown contracts
      invoices/             CSVs
      reports/              CSVs + PNG charts
      images/               PNG logos / scans
      README.md, notes/*.md text files (content search targets)

Everything is generated deterministically (fixed seed) so tests can assert on
specific filenames/rows. Used both by the e2e test and the sandbox feedback
loop doc.

Usage:
    python scripts/gen_network_dir_fixtures.py /path/to/dir [--contracts 40]
"""
from __future__ import annotations

import argparse
import csv
import os
import random
import zipfile
from pathlib import Path

import matplotlib
matplotlib.use("Agg")  # headless
import matplotlib.pyplot as plt
import numpy as np
from PIL import Image, ImageDraw

VENDORS = [
    "Acme Corp", "Globex", "Initech", "Umbrella", "Soylent", "Stark Industries",
    "Wayne Enterprises", "Wonka Industries", "Cyberdyne", "Tyrell Corp",
    "Nakatomi", "Gekko & Co", "Vandelay", "Hooli", "Pied Piper",
]
STATUSES = ["active", "expired", "pending", "terminated"]


def _seed():
    random.seed(1234)
    np.random.seed(1234)


def _write_csv(path: Path, header, rows):
    path.parent.mkdir(parents=True, exist_ok=True)
    with open(path, "w", newline="") as fh:
        w = csv.writer(fh)
        w.writerow(header)
        w.writerows(rows)


def gen_contracts(root: Path, n: int) -> int:
    count = 0
    for i in range(n):
        vendor = random.choice(VENDORS)
        slug = vendor.lower().replace(" ", "_").replace("&", "and")
        year = random.choice([2022, 2023, 2024, 2025])
        value = random.randint(10_000, 2_000_000)
        status = random.choice(STATUSES)
        rows = [
            [f"CLAUSE-{j+1}", random.choice(["payment", "termination", "liability",
             "confidentiality", "renewal", "indemnity"]),
             random.randint(1, 36), round(random.uniform(0, value / 10), 2)]
            for j in range(random.randint(4, 12))
        ]
        _write_csv(
            root / "contracts" / f"contract_{slug}_{year}_{i:03d}.csv",
            ["clause_id", "clause_type", "term_months", "penalty_usd"],
            rows,
        )
        # A markdown summary the way search_files content-grep would find it.
        (root / "contracts").mkdir(parents=True, exist_ok=True)
        (root / "contracts" / f"contract_{slug}_{year}_{i:03d}.md").write_text(
            f"# Master Services Agreement — {vendor} ({year})\n\n"
            f"Status: **{status}**\n\nTotal contract value: ${value:,}\n\n"
            f"This agreement between the Company and {vendor} governs the "
            f"provision of services. Governing law: Delaware.\n"
        )
        count += 2
    return count


def gen_invoices(root: Path, n: int) -> int:
    for i in range(n):
        vendor = random.choice(VENDORS)
        rows = [
            [f"INV-{i:04d}-{k}", vendor, f"2025-{random.randint(1,12):02d}-{random.randint(1,28):02d}",
             round(random.uniform(100, 50_000), 2), random.choice(["paid", "unpaid", "overdue"])]
            for k in range(random.randint(3, 8))
        ]
        _write_csv(
            root / "invoices" / f"invoices_{i:03d}.csv",
            ["invoice_id", "vendor", "date", "amount_usd", "status"],
            rows,
        )
    return n


def gen_reports_and_charts(root: Path, n: int) -> int:
    count = 0
    for i in range(n):
        months = ["Jan", "Feb", "Mar", "Apr", "May", "Jun"]
        revenue = np.random.randint(50, 500, size=len(months))
        _write_csv(
            root / "reports" / f"revenue_q{i+1}.csv",
            ["month", "revenue_k_usd"],
            list(zip(months, revenue.tolist())),
        )
        # Matching PNG chart.
        fig, ax = plt.subplots(figsize=(4, 3))
        ax.bar(months, revenue, color="#4C78A8")
        ax.set_title(f"Revenue Q{i+1}")
        (root / "reports").mkdir(parents=True, exist_ok=True)
        fig.savefig(root / "reports" / f"revenue_q{i+1}.png", dpi=80)
        plt.close(fig)
        count += 2
    return count


def gen_images(root: Path, n: int) -> int:
    (root / "images").mkdir(parents=True, exist_ok=True)
    for i in range(n):
        img = Image.new("RGB", (240, 120), color=(random.randint(0, 255),
                        random.randint(0, 255), random.randint(0, 255)))
        d = ImageDraw.Draw(img)
        d.rectangle([10, 10, 230, 110], outline=(255, 255, 255), width=3)
        d.text((30, 50), f"LOGO {i:02d}", fill=(255, 255, 255))
        img.save(root / "images" / f"logo_{i:02d}.png")
    return n


def _make_docx(path: Path, paragraphs) -> None:
    """Write a minimal but valid .docx (OOXML zip) — no python-docx needed."""
    body = "".join(
        f"<w:p><w:r><w:t>{p}</w:t></w:r></w:p>" for p in paragraphs
    )
    document = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        f"<w:body>{body}</w:body></w:document>"
    )
    content_types = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
        '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
        '<Default Extension="xml" ContentType="application/xml"/>'
        '<Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>'
        "</Types>"
    )
    rels = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>'
        "</Relationships>"
    )
    path.parent.mkdir(parents=True, exist_ok=True)
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("[Content_Types].xml", content_types)
        z.writestr("_rels/.rels", rels)
        z.writestr("word/document.xml", document)


def _make_pdf(path: Path, title: str, lines) -> None:
    """matplotlib PDF backend embeds selectable text, so pypdf/pdfminer can
    extract it — no reportlab needed."""
    path.parent.mkdir(parents=True, exist_ok=True)
    fig = plt.figure(figsize=(8.5, 11))
    fig.text(0.1, 0.9, title, fontsize=16, weight="bold")
    fig.text(0.1, 0.8, "\n".join(lines), fontsize=11, va="top")
    fig.savefig(path)
    plt.close(fig)


def gen_documents(root: Path, n: int) -> int:
    """A handful of contracts as pdf / docx / pptx with searchable content."""
    from pptx import Presentation
    from pptx.util import Inches

    count = 0
    for i in range(n):
        vendor = VENDORS[i % len(VENDORS)]
        slug = vendor.lower().replace(" ", "_").replace("&", "and")
        term = random.choice(["indemnity", "auto-renewal", "arbitration", "force majeure"])
        value = random.randint(50_000, 5_000_000)

        # PDF
        _make_pdf(
            root / "documents" / f"contract_{slug}_{i:02d}.pdf",
            f"Master Services Agreement — {vendor}",
            [f"Governing law: Delaware.", f"Total contract value: ${value:,}.",
             f"Key clause: {term}.", "Confidentiality applies to all parties."],
        )
        # DOCX
        _make_docx(
            root / "documents" / f"agreement_{slug}_{i:02d}.docx",
            [f"Services Agreement with {vendor}",
             f"This agreement includes a {term} clause.",
             f"Contract value: ${value:,}. Governing law: Delaware."],
        )
        # PPTX
        prs = Presentation()
        s = prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.title.text = f"{vendor} — Deal Review"
        tb = s.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3)).text_frame
        tb.text = f"Vendor: {vendor}\nValue: ${value:,}\nClause: {term}\nStatus: under review"
        (root / "documents").mkdir(parents=True, exist_ok=True)
        prs.save(root / "documents" / f"review_{slug}_{i:02d}.pptx")
        count += 3
    return count


def gen_spreadsheets(root: Path, n: int) -> int:
    """A few multi-sheet .xlsx workbooks (budgets) with searchable content."""
    import pandas as pd

    (root / "spreadsheets").mkdir(parents=True, exist_ok=True)
    for i in range(n):
        vendor = VENDORS[i % len(VENDORS)]
        budget = pd.DataFrame({
            "category": ["licenses", "services", "support", "training"],
            "vendor": [vendor] * 4,
            "amount_usd": [random.randint(1000, 90000) for _ in range(4)],
            "quarter": [f"Q{random.randint(1,4)}" for _ in range(4)],
        })
        headcount = pd.DataFrame({
            "team": ["eng", "sales", "ops"],
            "count": [random.randint(3, 40) for _ in range(3)],
        })
        path = root / "spreadsheets" / f"budget_{vendor.lower().replace(' ', '_')}_{i:02d}.xlsx"
        with pd.ExcelWriter(path) as xl:
            budget.to_excel(xl, sheet_name="budget", index=False)
            headcount.to_excel(xl, sheet_name="headcount", index=False)
    return n


def gen_notes(root: Path) -> int:
    (root / "notes").mkdir(parents=True, exist_ok=True)
    (root / "README.md").write_text(
        "# Vendor Document Store\n\nContracts, invoices, reports and images for "
        "the vendor management workflow. The Acme Corp master agreement is the "
        "reference template.\n"
    )
    (root / "notes" / "renewal_checklist.md").write_text(
        "# Renewal Checklist\n\n- Confirm Umbrella contract status\n"
        "- Review payment clauses\n- Flag any terminated agreements\n"
    )
    return 2


def generate(root: Path, contracts: int = 40, invoices: int = 20,
             reports: int = 4, images: int = 8, documents: int = 4,
             spreadsheets: int = 4) -> dict:
    _seed()
    root.mkdir(parents=True, exist_ok=True)
    stats = {
        "contracts": gen_contracts(root, contracts),
        "invoices": gen_invoices(root, invoices),
        "reports": gen_reports_and_charts(root, reports),
        "images": gen_images(root, images),
        "documents": gen_documents(root, documents),
        "spreadsheets": gen_spreadsheets(root, spreadsheets),
        "notes": gen_notes(root),
    }
    total = sum(1 for _ in root.rglob("*") if _.is_file())
    stats["total_files"] = total
    return stats


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("root", help="Target directory to populate")
    ap.add_argument("--contracts", type=int, default=40)
    ap.add_argument("--invoices", type=int, default=20)
    ap.add_argument("--reports", type=int, default=4)
    ap.add_argument("--images", type=int, default=8)
    ap.add_argument("--documents", type=int, default=4)
    ap.add_argument("--spreadsheets", type=int, default=4)
    args = ap.parse_args()
    stats = generate(Path(args.root).expanduser().resolve(),
                     args.contracts, args.invoices, args.reports, args.images,
                     args.documents, args.spreadsheets)
    print(f"Populated {args.root}:")
    for k, v in stats.items():
        print(f"  {k}: {v}")


if __name__ == "__main__":
    main()
